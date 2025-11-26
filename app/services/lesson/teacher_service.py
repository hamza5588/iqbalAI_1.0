"""
Teacher-focused lesson service for creating and managing lessons
"""
import os
import logging
import tempfile
from typing import Any, Dict, List, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
import base64
import threading

# Disable tqdm threading to prevent "cannot start new thread" errors
os.environ['TQDM_DISABLE'] = '1'
os.environ['TOKENIZERS_PARALLELISM'] = 'false'
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename
from langchain_core.documents import Document
from langchain_core.output_parsers import PydanticOutputParser
from langchain_core.prompts import PromptTemplate
from docx import Document as DocxDocument
from docx.shared import Inches
from io import BytesIO
from datetime import datetime
import json

# Try to import optional dependencies
try:
    import camelot
    CAMELOT_AVAILABLE = True
except ImportError:
    CAMELOT_AVAILABLE = False

try:
    from PIL import Image
    import fitz  # PyMuPDF for image extraction
    PIL_AVAILABLE = True
    PYMUPDF_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    PYMUPDF_AVAILABLE = False

from .base_service import BaseLessonService
from .models import LessonResponse, LessonPlan
from .rag_service import RAGService

from app.models.models import LessonModel

logger = logging.getLogger(__name__)

# Set up detailed teacher service logging
teacher_logger = logging.getLogger('teacher_service')
teacher_logger.setLevel(logging.INFO)

# Create teacher.log file handler
if not os.path.exists('logs'):
    os.makedirs('logs')

teacher_handler = logging.FileHandler('logs/lesson.log')
teacher_handler.setLevel(logging.INFO)

# Create formatter for teacher logs
teacher_formatter = logging.Formatter(
    '%(asctime)s - TEACHER_SERVICE - %(levelname)s - %(funcName)s:%(lineno)d - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
teacher_handler.setFormatter(teacher_formatter)
teacher_logger.addHandler(teacher_handler)

from typing import Literal
from pydantic import BaseModel

from langchain_core.prompts import PromptTemplate
from langchain_core.output_parsers import JsonOutputParser
# from langchain.memory import ConversationBufferMemory

from langchain_core.chat_history import BaseChatMessageHistory, InMemoryChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_core.chat_history import BaseChatMessageHistory, InMemoryChatMessageHistory
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.messages import HumanMessage, AIMessage
from langchain_community.vectorstores import FAISS
from langchain_core.runnables import RunnablePassthrough
from pydantic import BaseModel, Field
import re
from langchain_core.runnables import RunnableLambda
from langchain_core.runnables.config import RunnableConfig
from langchain_openai import ChatOpenAI



class lesson_response(BaseModel):
    complete_lesson: Literal["yes", "no"] = Field(
        ...,
        description="check if the lesson is generated or not if generated ('yes') or not ('no')."
    )

class InteractiveChatResponse(BaseModel):
    """Structured response for interactive chat"""
    ai_response: str = Field(..., description="The AI's response text")
    complete_lesson: Literal["yes", "no"] = Field(..., description="Whether the complete lesson has been generated ('yes') or still in draft/outline stage ('no')")

def check_lesson_response(text: str, groq_api_key: str):
    """Check if the AI response indicates a complete lesson has been generated"""
    # llm = ChatGroq(
    #     groq_api_key=groq_api_key,
    #     model_name="llama-3.1-8b-instant",
    #     temperature=0.1
    # )
    vllm_api_base = os.getenv('VLLM_API_BASE', 'http://69.28.92.113:8000/v1')
    vllm_model = os.getenv('VLLM_MODEL', 'Qwen/Qwen2-VL-2B-Instruct')
    llm = ChatOpenAI(
        openai_api_key="EMPTY",
        openai_api_base=vllm_api_base,
        model_name=vllm_model,
        temperature=0.7,
        max_tokens=512,
    )
    
    # Create a prompt to analyze if the response is a complete lesson or just an outline/draft
    analysis_prompt = f"""Analyze the following AI response and determine if it contains a COMPLETE LESSON or just an OUTLINE/DRAFT.

Complete Lesson indicators:
- Contains full lesson content with detailed explanations
- Includes actual lesson sections with comprehensive content
- Has detailed examples, code snippets, or explanations
- Ready to use for teaching (not just a structure/outline)

Outline/Draft indicators:
- Only shows structure (Learning Objectives, Main Topics list, etc.)
- Asks for confirmation before generating full lesson
- Contains phrases like "Does this outline work", "Should I generate", "Is this perfect now"
- Only shows headings/sections without full content

AI Response:
{text}

Determine if this is a complete lesson (yes) or still a draft/outline (no)."""

    llm_with_structured_output = llm.with_structured_output(lesson_response)
    return llm_with_structured_output.invoke(analysis_prompt)

class TeacherLessonService(BaseLessonService):
    """
    Teacher-focused lesson generation workflow:
    1. Teacher selects 'Generate Lesson' from the sidebar.
    2. Uploads a PDF file (e.g., textbook, handout).
    3. Teacher enters details about the lesson plan to be generated (title, objectives, focus areas, etc.).
    4. iqbalAI processes the file and user inputs, generates a structured lesson plan.
    5. Teacher can review, ask follow-up questions, edit/refine lesson parts.
    6. After final screening and edits, the final document is ready.
    7. User/teacher can download the final version.
    """
    
    # Class-level dictionary to persist chat histories across service instances
    # This ensures conversation history persists across HTTP requests
    _chat_histories = {}
    
    def __init__(self, groq_api_key: str):
        super().__init__(groq_api_key)
        self.rag_service = RAGService()
        self.lesson_vector_stores = {}  # Store vector DBs for each lesson
        teacher_logger.info("RAG service initialized")
        # Use class-level chat_histories to persist across instances
        teacher_logger.info("RAG service initialized with persistent memory")
        
        # Initialize separate multimodal LLM for image descriptions
        vllm_api_base = os.getenv('VLLM_API_BASE', 'http://69.28.92.113:8000/v1')
        vllm_multimodal_model = os.getenv('VLLM_MULTIMODAL_MODEL', 'Qwen/Qwen2-VL-2B-Instruct')
        vllm_timeout = int(os.getenv('VLLM_TIMEOUT', 600))
        
        self.multimodal_llm = ChatOpenAI(
            openai_api_key="EMPTY",
            openai_api_base=vllm_api_base,
            model_name=vllm_multimodal_model,
            temperature=0.7,
            max_tokens=1024,
            timeout=vllm_timeout,
        )
        teacher_logger.info(f"Multimodal LLM initialized for image descriptions: {vllm_multimodal_model}")
    
    def _detect_pages_with_tables(self, file_path: str) -> List[int]:
        """Quickly detect which pages likely contain tables by scanning for table-like structures"""
        pages_with_tables = []
        
        try:
            if not PYMUPDF_AVAILABLE:
                teacher_logger.warning("PyMuPDF not available for table detection, will process all pages")
                return []
            
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            teacher_logger.info(f"Scanning {len(doc)} pages for table detection...")
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Get drawings (lines) on the page - tables typically have many horizontal/vertical lines
                drawings = page.get_drawings()
                
                # Count horizontal and vertical lines (indicators of tables)
                horizontal_lines = 0
                vertical_lines = 0
                
                for drawing in drawings:
                    # Check if it's a line (has a path with 2 points)
                    if 'rect' in drawing:
                        rect = drawing['rect']
                        width = abs(rect.x1 - rect.x0)
                        height = abs(rect.y1 - rect.y0)
                        
                        # Horizontal line (width >> height)
                        if width > height * 3 and height < 5:
                            horizontal_lines += 1
                        # Vertical line (height >> width)
                        elif height > width * 3 and width < 5:
                            vertical_lines += 1
                
                # If we find multiple horizontal and vertical lines, likely a table
                # Threshold: at least 3 horizontal and 2 vertical lines suggests a table
                if horizontal_lines >= 3 and vertical_lines >= 2:
                    pages_with_tables.append(page_num + 1)  # Camelot uses 1-based page numbers
                    teacher_logger.info(f"Page {page_num + 1} likely contains tables (found {horizontal_lines} horizontal, {vertical_lines} vertical lines)")
            
            doc.close()
            
            if pages_with_tables:
                teacher_logger.info(f"Detected potential tables on pages: {pages_with_tables}")
            else:
                teacher_logger.info("No pages with table-like structures detected")
            
            return pages_with_tables
            
        except Exception as e:
            teacher_logger.warning(f"Error detecting pages with tables: {str(e)}. Will process all pages.")
            return []
    
    def _extract_tables_from_pdf(self, file_path: str, filename: str) -> str:
        """Extract tables from PDF using Camelot and return as Markdown.
        First detects which pages contain tables, then only processes those pages."""
        if not CAMELOT_AVAILABLE:
            teacher_logger.warning("Camelot not available, skipping table extraction")
            return ""
        
        if not filename.lower().endswith('.pdf'):
            teacher_logger.info(f"Table extraction only supported for PDF files, got {filename}")
            return ""
        
        try:
            teacher_logger.info(f"Extracting tables from {filename}")
            
            # Step 1: Detect which pages contain tables
            pages_with_tables = self._detect_pages_with_tables(file_path)
            
            if not pages_with_tables:
                teacher_logger.info("No pages with tables detected, skipping table extraction")
                return ""
            
            # Step 2: Only process pages that likely contain tables
            # Convert list to comma-separated string for Camelot (e.g., "1,3,5")
            pages_str = ','.join(map(str, pages_with_tables))
            teacher_logger.info(f"Processing only pages with detected tables: {pages_str}")
            
            # Use camelot to extract tables from specific pages only
            tables = camelot.read_pdf(file_path, pages=pages_str, flavor='lattice')
            
            if len(tables) == 0:
                teacher_logger.info("No tables found in detected pages")
                return ""
            
            teacher_logger.info(f"Found {len(tables)} tables in pages {pages_str}")
            
            # Convert tables to markdown format
            markdown_tables = []
            for i, table in enumerate(tables, 1):
                df = table.df
                # Include page number in table header if available
                page_info = f" (Page {table.page})" if hasattr(table, 'page') else ""
                markdown_tables.append(f"\n### Table {i}{page_info}\n\n")
                # Convert DataFrame to markdown table
                markdown_table = df.to_markdown(index=False)
                markdown_tables.append(markdown_table)
                markdown_tables.append("\n")
            
            result = "\n".join(markdown_tables)
            teacher_logger.info(f"Extracted {len(tables)} tables as markdown from {len(pages_with_tables)} pages")
            return result
            
        except Exception as e:
            teacher_logger.error(f"Error extracting tables: {str(e)}")
            return ""
    
    def _extract_images_and_describe(self, file_path: str, filename: str) -> str:
        """Extract images from PDF and generate descriptions using LLM in parallel batches"""
        if not PYMUPDF_AVAILABLE:
            teacher_logger.warning("PyMuPDF not available, skipping image extraction")
            return ""
        
        if not filename.lower().endswith('.pdf'):
            teacher_logger.info(f"Image extraction only supported for PDF files, got {filename}")
            return ""
        
        try:
            teacher_logger.info(f"Extracting images from {filename}")
            import fitz  # PyMuPDF
            
            doc = fitz.open(file_path)
            
            # First, collect all images with their metadata
            image_tasks = []
            for page_num in range(len(doc)):
                page = doc[page_num]
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    try:
                        # Get image data
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]
                        
                        # Convert to base64 for LLM
                        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                        
                        # Get surrounding text context from the page (if available)
                        page_text = page.get_text("text")
                        context_snippet = ""
                        if page_text:
                            # Get first 200 characters of page text as context
                            context_snippet = page_text[:200].strip().replace('\n', ' ')
                        
                        # Store image task for parallel processing
                        image_tasks.append({
                            'base64': image_base64,
                            'ext': image_ext,
                            'page_num': page_num + 1,
                            'img_index': img_index + 1,
                            'context': context_snippet
                        })
                    except Exception as e:
                        teacher_logger.warning(f"Error preparing image {img_index} on page {page_num}: {str(e)}")
                        continue
            
            doc.close()
            
            if not image_tasks:
                teacher_logger.info("No images found in PDF")
                return ""
            
            teacher_logger.info(f"Found {len(image_tasks)} images, processing in parallel batches of 4")
            
            # Process images in parallel batches of 4
            image_descriptions = []
            batch_size = 4
            
            for batch_start in range(0, len(image_tasks), batch_size):
                batch = image_tasks[batch_start:batch_start + batch_size]
                teacher_logger.info(f"Processing image batch {batch_start // batch_size + 1} ({len(batch)} images)")
                
                # Process batch in parallel
                with ThreadPoolExecutor(max_workers=batch_size) as executor:
                    batch_futures = {
                        executor.submit(
                            self._process_single_image,
                            task['base64'],
                            task['ext'],
                            task['page_num'],
                            task['img_index'],
                            task['context']
                        ): task for task in batch
                    }
                    
                    # Collect results from batch
                    for future in as_completed(batch_futures):
                        task = batch_futures[future]
                        try:
                            description_entry = future.result()
                            if description_entry:
                                image_descriptions.append(description_entry)
                        except Exception as e:
                            teacher_logger.warning(f"Error processing image {task['page_num']}-{task['img_index']}: {str(e)}")
            
            if not image_descriptions:
                teacher_logger.info("No image descriptions generated")
                return ""
            
            result = "\n".join(image_descriptions)
            teacher_logger.info(f"Extracted and described {len(image_descriptions)} images")
            return result
            
        except Exception as e:
            teacher_logger.error(f"Error extracting images: {str(e)}")
            return ""
    
    def _process_single_image(self, image_base64: str, image_ext: str, page_num: int, img_index: int, context: str) -> str:
        """Process a single image and return formatted description entry"""
        try:
            # Generate description using LLM with context
            description = self._describe_image_with_llm(
                image_base64, 
                image_ext, 
                page_num=page_num,
                image_index=img_index,
                context=context
            )
            
            if not description:
                return ""
            
            # Extract image type from description for better searchability
            image_type_keywords = ""
            description_lower = description.lower()
            if any(word in description_lower for word in ['pie chart', 'piechart']):
                image_type_keywords = "pie chart, chart, graph, visualization, data visualization"
            elif any(word in description_lower for word in ['bar chart', 'bar graph', 'barchart']):
                image_type_keywords = "bar chart, bar graph, chart, graph, visualization, data visualization"
            elif any(word in description_lower for word in ['line graph', 'line chart']):
                image_type_keywords = "line graph, line chart, chart, graph, visualization, data visualization"
            elif any(word in description_lower for word in ['diagram', 'flowchart']):
                image_type_keywords = "diagram, flowchart, illustration, visual"
            elif any(word in description_lower for word in ['table', 'data table']):
                image_type_keywords = "table, data table, information table"
            else:
                image_type_keywords = "image, picture, visual, illustration"
            
            # Create a comprehensive image entry with context and searchable keywords
            image_entry = f"""
### Image on Page {page_num} (Image {img_index})

**Location:** Page {page_num} of the document
**Image Type:** {image_type_keywords}

**Description:**
{description}

**Context:** This image appears on page {page_num} of the document. {f"Surrounding context: {context}" if context else ""}

**Searchable Keywords:** image, picture, visual, page {page_num}, {image_type_keywords}
"""
            teacher_logger.info(f"Generated description for image {page_num}-{img_index}")
            return image_entry
            
        except Exception as e:
            teacher_logger.warning(f"Error processing image {page_num}-{img_index}: {str(e)}")
            return ""
    
    def _describe_image_with_llm(self, image_base64: str, image_ext: str, page_num: int = None, image_index: int = None, context: str = "") -> str:
        """Generate description of image using multimodal LLM with enhanced context and retry logic"""
        import time
        max_retries = 3
        retry_delay = 2  # seconds
        
        # Check image size - if too large, it might cause connection issues
        image_size_mb = len(image_base64) * 3 / 4 / 1024 / 1024  # Approximate size in MB
        if image_size_mb > 10:  # If image is larger than 10MB, skip description
            teacher_logger.warning(f"Image too large ({image_size_mb:.2f}MB), skipping description")
            return f"[Image extracted from document on page {page_num if page_num else 'unknown'} - image too large for description]"
        
        for attempt in range(max_retries):
            try:
                from langchain_core.messages import HumanMessage
                
                # Build context information
                context_info = ""
                if page_num:
                    context_info += f"This image appears on page {page_num} of the document. "
                if context:
                    context_info += f"The surrounding text context: {context[:150]}... "
                
                # Create enhanced prompt for image description that includes context references
                prompt_text = f"""Please analyze and describe the provided image in thorough detail using the structured format below. Your response will be used for RAG, so ensure the description is clear, explicit, and information-dense.

1. *Image Type & Overall Message*

Begin by clearly identifying the type of image using phrases such as:

“This is a pie chart showing…”

“This is a bar graph displaying…”

“This line graph indicates…”

“This diagram illustrates…”

“This image shows…” or “This picture depicts…”

Then explicitly state what the image is telling us.

2. *Data & Information (for charts/graphs/tables)*

If the image contains data visualization, clearly describe:

What data is being presented

What the chart/graph is telling us

Key values, percentages, or labels

Trends, increases, decreases, patterns

Main insights or conclusions

Use phrases such as:

“The chart is telling us that…”

“This graph shows that…”

“The data indicates that…”

3. *Visual Elements*

Describe all visible elements, including:

Colors

Labels, legends, axes, and scales

Titles, subtitles, annotations

Structure, layout, shapes, icons

4. *What Is Happening in the Image* (for scenes, photos, illustrations)

If the image shows a scene or action, describe:

People, objects, or entities present

Actions taking place

Environment or setting

Notable activities or sequences

Use phrases such as:

“What is happening in this image: …”

“This image shows the following: …”

5. *Context & Purpose*

Explain the overall purpose of the image:

What idea or concept it illustrates

Why this image or chart is useful

What educational or informational purpose it serves

The main story or insight communicated

6. *Answer Format Guide*

Structure your descriptions so they naturally answer common question forms:

“What is the pie chart telling?” → “This pie chart is telling us that…”

“What does this graph show?” → “This graph shows that…”

“What data is presented?” → “The data presented includes…”

“What is happening in this image?” → “What is happening in this image: …”

{context_info}

Provide a detailed, comprehensive description that explicitly answers questions about what the image is telling, showing, or what is happening in it. Use natural, conversational language that makes it easy to understand and query the image content."""
                
                # Format image extension (remove leading dot if present)
                img_ext = image_ext.lstrip('.') if image_ext else 'png'
                
                # Create multimodal message with text and image
                message = HumanMessage(
                    content=[
                        {"type": "text", "text": prompt_text},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/{img_ext};base64,{image_base64}"
                            }
                        }
                    ]
                )
                
                teacher_logger.info(f"Requesting image description from multimodal LLM (attempt {attempt + 1}/{max_retries}, format: {img_ext}, page: {page_num}, size: {image_size_mb:.2f}MB)")
                
                # Invoke multimodal LLM (separate instance for image descriptions)
                response = self.multimodal_llm.invoke([message])
                
                if hasattr(response, 'content'):
                    description = response.content.strip()
                    teacher_logger.info(f"Generated image description: {len(description)} characters")
                    return description
                else:
                    description = str(response).strip()
                    teacher_logger.info(f"Generated image description: {len(description)} characters")
                    return description
                    
            except Exception as e:
                error_str = str(e)
                
                # Check if it's a "not multimodal model" error - don't retry for this
                is_not_multimodal_error = 'not a multimodal model' in error_str.lower() or 'not multimodal' in error_str.lower()
                if is_not_multimodal_error:
                    teacher_logger.error(f"Model does not support multimodal requests: {error_str[:200]}")
                    fallback_msg = f"[Image extracted from document on page {page_num if page_num else 'unknown'}"
                    if context:
                        fallback_msg += f" - Context: {context[:100]}..."
                    fallback_msg += " - Image description unavailable: model does not support multimodal requests]"
                    return fallback_msg
                
                # Check for connection errors that can be retried
                is_connection_error = any(keyword in error_str.lower() for keyword in [
                    'connection', 'refused', 'timeout', 'connect', 'network', '10061'
                ])
                
                if attempt < max_retries - 1 and is_connection_error:
                    wait_time = retry_delay * (attempt + 1)  # Exponential backoff
                    teacher_logger.warning(f"Connection error on attempt {attempt + 1}/{max_retries}: {error_str[:200]}. Retrying in {wait_time}s...")
                    time.sleep(wait_time)
                    continue
                else:
                    teacher_logger.error(f"Error generating image description (attempt {attempt + 1}/{max_retries}): {error_str}", exc_info=True)
                    # Return a fallback message that still provides some value
                    fallback_msg = f"[Image extracted from document on page {page_num if page_num else 'unknown'}"
                    if context:
                        fallback_msg += f" - Context: {context[:100]}..."
                    fallback_msg += " - Description generation unavailable due to error]"
                    return fallback_msg
        
        # If all retries failed
        return f"[Image extracted from document on page {page_num if page_num else 'unknown'} - description generation failed after {max_retries} attempts]"

    def process_file(self, file: FileStorage, lesson_details: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
        """
        Process an uploaded file and create vector DB for RAG.
        Supports parallel extraction of text, tables, and images.
        Returns a greeting message instead of generating lesson immediately.
        Lesson generation happens only when user sends a query via interactive_chat.
        """
        teacher_logger.info(f"=== TEACHER FILE PROCESSING STARTED ===")
        teacher_logger.info(f"File: {file.filename if file else 'None'}")
        teacher_logger.info(f"Lesson details: {json.dumps(lesson_details, indent=2) if lesson_details else 'None'}")
        
        # Get extraction flags
        table_extraction = lesson_details.get('table_extraction', False) if lesson_details else False
        image_extraction = lesson_details.get('image_extraction', False) if lesson_details else False
        
        temp_path = None
        try:
            if not file or not file.filename:
                teacher_logger.warning("No file provided")
                return {"error": "No file provided"}
            if not self.allowed_file(file.filename):
                teacher_logger.warning(f"Unsupported file type: {file.filename}")
                return {"error": "File type not supported. Please upload PDF, DOC, DOCX, or TXT files."}
            
            teacher_logger.info(f"File validation passed: {file.filename}")
            teacher_logger.info(f"Table extraction: {table_extraction}, Image extraction: {image_extraction}")
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{secure_filename(file.filename)}") as temp_file:
                temp_path = temp_file.name
                file.save(temp_path)
            
            teacher_logger.info(f"File saved to temporary path: {temp_path}")
            
            # Step 1: Process text and tables first (synchronous)
            text_content = ""
            table_content = ""
            
            with ThreadPoolExecutor(max_workers=2) as executor:
                futures = {}
                
                # Always extract text
                futures['text'] = executor.submit(self._load_document, temp_path, file.filename)
                
                # Extract tables if enabled
                if table_extraction:
                    futures['tables'] = executor.submit(self._extract_tables_from_pdf, temp_path, file.filename)
                
                # Wait for text and tables to complete
                for key, future in futures.items():
                    try:
                        result = future.result()
                        if key == 'text':
                            documents = result
                            if documents:
                                text_content = "\n".join([doc.page_content for doc in documents])
                        elif key == 'tables':
                            table_content = result if result else ""
                    except Exception as e:
                        teacher_logger.error(f"Error in {key} extraction: {str(e)}")
                        if key == 'text':
                            return {"error": f"Could not extract content from the file: {str(e)}"}
            
            if not text_content.strip():
                teacher_logger.error("No readable content found in the file")
                return {"error": "No readable content found in the file"}
            
            # Step 2: Create initial content with text and tables
            initial_content = text_content
            if table_content:
                initial_content += "\n\n## Extracted Tables\n" + table_content
                teacher_logger.info(f"Added {len(table_content)} characters of table content")
            
            teacher_logger.info(f"Initial content extracted: {len(initial_content)} characters (Text: {len(text_content)}, Tables: {len(table_content)})")
            
            # Step 3: Create initial vector DB with text and tables
            initial_documents = [Document(page_content=initial_content, metadata={"source": file.filename})]
            rag_result = self.rag_service.process_document(initial_documents, file.filename)
            if 'error' in rag_result:
                teacher_logger.error(f"RAG processing failed: {rag_result['error']}")
                return rag_result
            
            # Store the original document's RAG service
            if rag_result['use_rag']:
                self._store_original_document_rag(file.filename)
                teacher_logger.info("Initial vector DB created with text and tables")
            
            # Step 4: If image extraction is enabled, process in background
            images_processing = False
            if image_extraction:
                images_processing = True
                teacher_logger.info("Starting background image processing...")
                
                # Start background thread for image processing
                def process_images_background():
                    try:
                        teacher_logger.info("Background image processing started")
                        image_content = self._extract_images_and_describe(temp_path, file.filename)
                        
                        if image_content:
                            # Update vector DB with image descriptions
                            updated_content = initial_content + "\n\n## Image Descriptions\n" + image_content
                            updated_documents = [Document(page_content=updated_content, metadata={"source": file.filename})]
                            
                            # Re-process with updated content
                            self.rag_service.process_document(updated_documents, file.filename)
                            teacher_logger.info(f"Vector DB updated with image descriptions: {len(image_content)} characters")
                        else:
                            teacher_logger.info("No image descriptions generated")
                    except Exception as e:
                        teacher_logger.error(f"Error in background image processing: {str(e)}", exc_info=True)
                    finally:
                        # Clean up temp file in background thread
                        if temp_path and os.path.exists(temp_path):
                            try:
                                os.remove(temp_path)
                                teacher_logger.info(f"Temporary file cleaned up in background: {temp_path}")
                            except Exception as e:
                                teacher_logger.warning(f"Could not remove temporary file {temp_path}: {str(e)}")
                
                # Start background thread (daemon=True so it doesn't block app shutdown)
                bg_thread = threading.Thread(target=process_images_background, daemon=True)
                bg_thread.start()
                teacher_logger.info("Background image processing thread started")
            else:
                # Clean up temp file immediately if no image processing
                if temp_path and os.path.exists(temp_path):
                    try:
                        os.remove(temp_path)
                        teacher_logger.info(f"Temporary file cleaned up: {temp_path}")
                    except Exception as e:
                        teacher_logger.warning(f"Could not remove temporary file {temp_path}: {str(e)}")
            
            teacher_logger.info("Vector DB created successfully. Skipping LLM call - will respond on user query.")
            
            # Return greeting message with image processing status
            greeting_message = f"Hello! I'm Prof. Potter, here to help you prepare your lesson plan. Your file '{file.filename}' has been uploaded and processed successfully. I've analyzed the document and I'm ready to help you create a lesson from this content. How would you like me to help you create a lesson?"
            
            teacher_logger.info("=== TEACHER FILE PROCESSING COMPLETED (Text and tables done, images in background) ===")
            
            return {
                "lesson": greeting_message,
                "docx_bytes": None,  # No DOCX generated at upload time
                "filename": None,
                "file_processed": True,
                "filename_processed": file.filename,
                "images_processing": images_processing  # Flag to indicate images are processing
            }
        except Exception as e:
            teacher_logger.error(f"File processing failed: {str(e)}")
            logger.error(f"Error processing file: {str(e)}", exc_info=True)
            return {
                "error": "Failed to process file",
                "details": str(e)
            }
        finally:
            if temp_path and os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                    teacher_logger.info(f"Temporary file cleaned up: {temp_path}")
                except Exception as e:
                    teacher_logger.warning(f"Could not remove temporary file {temp_path}: {str(e)}")
                    logger.warning(f"Could not remove temporary file {temp_path}: {str(e)}")

    
    
    def get_session_history(self, session_id: str) -> BaseChatMessageHistory:
        """Get or create chat history for a session"""
        # Use class-level dictionary to persist history across service instances
        if session_id not in TeacherLessonService._chat_histories:
            TeacherLessonService._chat_histories[session_id] = InMemoryChatMessageHistory()
            teacher_logger.info(f"Created new chat history for session: {session_id}")
        else:
            history = TeacherLessonService._chat_histories[session_id]
            message_count = len(history.messages) if hasattr(history, 'messages') else 0
            teacher_logger.info(f"Retrieved existing chat history for session: {session_id} ({message_count} messages)")
        return TeacherLessonService._chat_histories[session_id]

    def _estimate_tokens(self, text: str) -> int:
        """Rough token estimation: ~4 characters per token"""
        return len(text) // 4
    
    def _truncate_text(self, text: str, max_tokens: int) -> str:
        """Truncate text to fit within token limit"""
        max_chars = max_tokens * 4
        if len(text) <= max_chars:
            return text
        # Truncate and add indicator
        return text[:max_chars - 50] + "... [truncated]"
    
    def format_context(self, relevant_chunks: List[Document], max_tokens: int = 2000) -> str:
        """Format context and escape curly braces to prevent LangChain template errors"""
        rag_context = "\n\n".join([doc.page_content for doc in relevant_chunks])
        # Truncate if too long
        rag_context = self._truncate_text(rag_context, max_tokens)
        # Escape curly braces in context content (this is a value, not a template string)
        # All braces in the context value should be escaped since it will be inserted into {context}
        escaped_context = rag_context.replace("{", "{{").replace("}", "}}")
        return escaped_context
    
    def _format_context_for_system_prompt(self, relevant_chunks: List[Document]) -> str:
        """Format context for system prompt without escaping (will be escaped later)"""
        return "\n\n".join([doc.page_content for doc in relevant_chunks])

    def interactive_chat(
    self, 
    lesson_id: int, 
    user_query: str, 
    session_id: str = None, 
    subject: str = None, 
    grade_level: str = None, 
    focus_area: str = None, 
    document_uploaded: bool = False, 
    document_filename: str = None
) -> InteractiveChatResponse:
        """Interactive chat with Prof. Potter for lesson creation"""
        
        teacher_logger.info("=== INTERACTIVE CHAT STARTED ===")
        
        # Use lesson_id as session_id
        if not session_id:
            session_id = f"lesson_{lesson_id}"
        
        try:
            # Step 1: Load vector DB
            vector_db = FAISS.load_local(
                "vector_store.faiss", 
                self.rag_service.embeddings, 
                allow_dangerous_deserialization=True
            )
            retriever = vector_db.as_retriever(search_type="similarity", search_kwargs={"k": 10})
            teacher_logger.info("Vector DB loaded successfully")
            
            # Step 2: Store form data
            form_context = {
                'subject': subject or focus_area,
                'grade_level': grade_level,
                'document_uploaded': document_uploaded,
                'document_filename': document_filename
            }
            
            # Step 3: Retrieve relevant context from vector store using user query
            docs = retriever.invoke(user_query)
            context = "\n\n".join([doc.page_content for doc in docs])
            teacher_logger.info(f"Retrieved {len(docs)} documents from vector store")
            
            # Step 4: Get chat history
            chat_history = self.get_session_history(session_id)
            teacher_logger.info(f"Chat history retrieved: {len(chat_history.messages) if hasattr(chat_history, 'messages') else 0} messages")
            
            # Step 5: Build system prompt with context
            base_system_prompt = self._get_system_prompt(rag_context=context, form_context=form_context)
            
            # Step 6: Build messages array
            from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
            langchain_messages = []
            
            # Add system message
            langchain_messages.append(SystemMessage(content=base_system_prompt))
            
            # Add chat history messages
            if hasattr(chat_history, 'messages'):
                for msg in chat_history.messages:
                    if hasattr(msg, 'type'):
                        if msg.type == "human":
                            langchain_messages.append(HumanMessage(content=msg.content if hasattr(msg, 'content') else str(msg)))
                        elif msg.type == "ai":
                            langchain_messages.append(AIMessage(content=msg.content if hasattr(msg, 'content') else str(msg)))
            
            # Add current user query
            langchain_messages.append(HumanMessage(content=user_query))
            
            # Step 7: Call LLM directly
            teacher_logger.info("Calling LLM...")
            response = self.llm.invoke(langchain_messages)
            
            # Extract response content
            if hasattr(response, 'content'):
                response_text = response.content
            elif isinstance(response, str):
                response_text = response
            else:
                response_text = str(response)
            
            teacher_logger.info(f"LLM response received: {len(response_text)} characters")
            
            # Step 8: Update chat history
            chat_history.add_message(HumanMessage(content=user_query))
            chat_history.add_message(AIMessage(content=response_text))
            teacher_logger.info(f"Chat history updated for session: {session_id}")
            
        except Exception as e:
            teacher_logger.error(f"Interactive chat error: {str(e)}", exc_info=True)
            raise
        
        # Step 9: Check if complete lesson generated
        try:
            lesson_check = check_lesson_response(response_text, self.api_key)
            complete_lesson_status = lesson_check.complete_lesson
            teacher_logger.info(f"Lesson completion check: {complete_lesson_status}")
        except Exception as e:
            teacher_logger.warning(f"Error checking lesson completion: {str(e)}")
            complete_lesson_status = "no"
        
        teacher_logger.info("=== INTERACTIVE CHAT COMPLETED ===")
        
        return InteractiveChatResponse(
            ai_response=response_text,
            complete_lesson=complete_lesson_status
        )
    
    def interactive_chat_stream(
        self, 
        lesson_id: int, 
        user_query: str, 
        session_id: str = None, 
        subject: str = None, 
        grade_level: str = None, 
        focus_area: str = None, 
        document_uploaded: bool = False, 
        document_filename: str = None
    ):
        """
        Interactive chat with streaming response for Prof. Potter.
        Yields chunks of the response as they are generated.
        Returns a generator that yields (chunk_text, is_complete, complete_lesson_status) tuples.
        """
        teacher_logger.info("=== INTERACTIVE CHAT STREAMING STARTED ===")
        
        # Use lesson_id as session_id
        if not session_id:
            session_id = f"lesson_{lesson_id}"
        
        response_text = ""
        complete_lesson_status = "no"
        
        try:
            # Step 1: Load vector DB
            vector_db = FAISS.load_local(
                "vector_store.faiss", 
                self.rag_service.embeddings, 
                allow_dangerous_deserialization=True
            )
            retriever = vector_db.as_retriever(search_type="similarity", search_kwargs={"k": 5})
            teacher_logger.info("Vector DB loaded successfully")
            
            # Step 2: Handle uploaded document content
            uploaded_doc_content = ""
            if document_uploaded and document_filename:
                try:
                    uploaded_doc_content = f"\n\n### Uploaded Document: {document_filename}\n[Document content]"
                    teacher_logger.info(f"Retrieved uploaded document: {document_filename}")
                except Exception as e:
                    teacher_logger.warning(f"Could not retrieve uploaded document: {str(e)}")
            
            # Step 3: Store form data
            form_context = {
                'subject': subject or focus_area,
                'grade_level': grade_level,
                'document_uploaded': document_uploaded,
                'document_filename': document_filename,
                'uploaded_content': uploaded_doc_content
            }
            
            # Step 4: Build system prompt
            base_system_prompt = self._get_system_prompt(form_context)
            teacher_logger.info("System prompt built")
            
            # Step 5: Get chat history
            chat_history = self.get_session_history(session_id)
            is_first_message = len(chat_history.messages) == 0 if hasattr(chat_history, 'messages') else True
            teacher_logger.info(f"Chat history retrieved: {len(chat_history.messages) if hasattr(chat_history, 'messages') else 0} messages")
            
            # Step 6: Enhance query if first message with document
            enhanced_query = user_query
            if is_first_message and document_uploaded:
                try:
                    overview_query = "What is this document about? Provide a brief summary."
                    overview_docs = retriever.invoke(overview_query)
                    if overview_docs:
                        doc_summary = "\n".join([doc.page_content[:200] for doc in overview_docs[:3]])
                        enhanced_query = f"{user_query}\n\n[Document Context: {doc_summary}...]"
                        teacher_logger.info("Query enhanced with document context")
                except Exception as e:
                    teacher_logger.warning(f"Could not retrieve document overview: {str(e)}")
            
            # Step 7: Retrieve context from vector store (reduced for faster CPU inference)
            docs = retriever.invoke(enhanced_query)
            context = self.format_context(docs, max_tokens=20000)  # Reduced from 1500 for faster processing
            teacher_logger.info(f"Retrieved {len(docs)} documents from vector store")
            
            # Step 8: Build messages array manually with token management
            messages = []
            
            # Add system message with context
            system_content = f"{base_system_prompt}\n\n### Knowledge Base Context:\n{context}{uploaded_doc_content}"
            messages.append({"role": "system", "content": system_content})
            
            # Add chat history messages (limit to last 5 messages for faster processing)
            if hasattr(chat_history, 'messages'):
                history_messages = list(chat_history.messages)
                if len(history_messages) > 5:  # Reduced from 10 to 5
                    history_messages = history_messages[-5:]
                
                for msg in history_messages:
                    if hasattr(msg, 'type'):
                        role = "user" if msg.type == "human" else "assistant"
                        content = msg.content if hasattr(msg, 'content') else str(msg)
                        if self._estimate_tokens(content) > 300:  # Reduced from 500 to 300
                            content = self._truncate_text(content, 300)
                        messages.append({"role": role, "content": content})
            
            # Add current user query
            messages.append({"role": "user", "content": enhanced_query})
            
            # Estimate total tokens before sending
            total_text = "\n".join([msg.get("content", "") for msg in messages])
            estimated_tokens = self._estimate_tokens(total_text)
            teacher_logger.info(f"Built message array with {len(messages)} messages, estimated tokens: {estimated_tokens}")
            
            # If estimated tokens exceed limit, reduce context further (lowered threshold for CPU)
            if estimated_tokens > 3000:  # Reduced from 5500 to 3000 for faster CPU inference
                teacher_logger.warning(f"Estimated tokens ({estimated_tokens}) exceed safe limit, reducing context")
                context = self.format_context(docs, max_tokens=500)  # Reduced from 800
                system_content = f"{base_system_prompt}\n\n### Knowledge Base Context:\n{context}{uploaded_doc_content}"
                messages[0] = {"role": "system", "content": system_content}
            
            # Step 9: Stream LLM response
            teacher_logger.info("Starting LLM streaming...")
            try:
                # Convert messages to LangChain message format
                from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
                langchain_messages = []
                for msg in messages:
                    role = msg.get("role", "user")
                    content = msg.get("content", "")
                    if role == "system":
                        langchain_messages.append(SystemMessage(content=content))
                    elif role == "user":
                        langchain_messages.append(HumanMessage(content=content))
                    elif role == "assistant":
                        langchain_messages.append(AIMessage(content=content))
                
                # Stream the response
                for chunk in self.llm.stream(langchain_messages):
                    if hasattr(chunk, 'content'):
                        chunk_text = chunk.content
                    else:
                        chunk_text = str(chunk)
                    
                    if chunk_text:
                        response_text += chunk_text
                        # Yield each chunk with is_complete=False
                        yield (chunk_text, False, "no")
                
                teacher_logger.info(f"LLM streaming completed: {len(response_text)} characters")
                
            except Exception as e:
                error_str = str(e)
                is_token_error = (
                    "413" in error_str or 
                    "too large" in error_str.lower() or 
                    "tokens per minute" in error_str.lower() or
                    "rate_limit_exceeded" in error_str.lower() or
                    "request too large" in error_str.lower()
                )
                
                if is_token_error:
                    teacher_logger.warning(f"Request too large, retrying with reduced context. Error: {error_str[:200]}")
                    # Retry with minimal context
                    context = self.format_context(docs, max_tokens=500)
                    system_content = f"{base_system_prompt}\n\n### Knowledge Base Context:\n{context}"
                    messages[0] = {"role": "system", "content": system_content}
                    
                    if hasattr(chat_history, 'messages'):
                        history_messages = list(chat_history.messages)
                        if len(history_messages) > 4:
                            history_messages = history_messages[-4:]
                        messages = [{"role": "system", "content": system_content}]
                        for msg in history_messages:
                            if hasattr(msg, 'type'):
                                role = "user" if msg.type == "human" else "assistant"
                                content = msg.content if hasattr(msg, 'content') else str(msg)
                                content = self._truncate_text(content, 300)
                                messages.append({"role": role, "content": content})
                        messages.append({"role": "user", "content": enhanced_query})
                    
                    # Retry streaming with LangChain message format
                    from langchain_core.messages import SystemMessage, HumanMessage, AIMessage
                    langchain_messages_retry = []
                    for msg in messages:
                        role = msg.get("role", "user")
                        content = msg.get("content", "")
                        if role == "system":
                            langchain_messages_retry.append(SystemMessage(content=content))
                        elif role == "user":
                            langchain_messages_retry.append(HumanMessage(content=content))
                        elif role == "assistant":
                            langchain_messages_retry.append(AIMessage(content=content))
                    
                    response_text = ""
                    for chunk in self.llm.stream(langchain_messages_retry):
                        if hasattr(chunk, 'content'):
                            chunk_text = chunk.content
                        else:
                            chunk_text = str(chunk)
                        
                        if chunk_text:
                            response_text += chunk_text
                            yield (chunk_text, False, "no")
                else:
                    # Re-raise if it's a different error
                    error_msg = f"\n\n[Error: {str(e)}]"
                    response_text += error_msg
                    yield (error_msg, False, "no")
                    raise
            
            # Step 10: Update chat history manually
            from langchain_core.messages import HumanMessage, AIMessage
            chat_history.add_message(HumanMessage(content=enhanced_query))
            chat_history.add_message(AIMessage(content=response_text))
            teacher_logger.info(f"Chat history updated for session: {session_id}")
            
        except Exception as e:
            teacher_logger.error(f"Interactive chat streaming error: {str(e)}", exc_info=True)
            # Yield error message
            error_msg = f"\n\n[Error: {str(e)}]"
            yield (error_msg, True, "no")
            return
        
        # Step 11: Check if complete lesson generated (after streaming completes)
        try:
            lesson_check = check_lesson_response(response_text, self.api_key)
            complete_lesson_status = lesson_check.complete_lesson
            teacher_logger.info(f"Lesson completion check: {complete_lesson_status}")
        except Exception as e:
            teacher_logger.warning(f"Error checking lesson completion: {str(e)}")
            complete_lesson_status = "no"
        
        # Step 12: Force cleanup
        import gc
        gc.collect()
        
        teacher_logger.info("=== INTERACTIVE CHAT STREAMING COMPLETED ===")
        
        # Yield final chunk with completion status
        yield ("", True, complete_lesson_status)
    
    # def interactive_chat(
    #     self, 
    #     lesson_id: int, 
    #     user_query: str, 
    #     session_id: str = None,
    #     subject: str = None,
    #     grade_level: str = None,
    #     focus_area: str = None,
    #     document_uploaded: bool = False,
    #     document_filename: str = None
    # ) -> InteractiveChatResponse:
    #     """Interactive chat with Prof. Potter for lesson creation"""
        
    #     # Use lesson_id as session_id to maintain context per lesson
    #     if not session_id:
    #         session_id = f"lesson_{lesson_id}"
        
    #     # Step 1: Load vector DB
    #     vector_db = FAISS.load_local(
    #         "vector_store.faiss",
    #         self.rag_service.embeddings,
    #         allow_dangerous_deserialization=True
    #     )
    #     retriever = vector_db.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        
    #     # Step 2: Handle uploaded document content
    #     uploaded_doc_content = ""
    #     if document_uploaded and document_filename:
    #         try:
    #             # TODO: Implement actual document retrieval from your storage
    #             # Example: uploaded_doc_content = self.get_uploaded_document(lesson_id, document_filename)
    #             uploaded_doc_content = f"\n\n### Uploaded Document: {document_filename}\n[Document content should be retrieved and inserted here]"
    #             teacher_logger.info(f"Retrieved uploaded document: {document_filename}")
    #         except Exception as e:
    #             teacher_logger.warning(f"Could not retrieve uploaded document: {str(e)}")
        
    #     # Step 3: Store form data in session context
    #     form_context = {
    #         'subject': subject or focus_area,
    #         'grade_level': grade_level,
    #         'document_uploaded': document_uploaded,
    #         'document_filename': document_filename,
    #         'uploaded_content': uploaded_doc_content  # Add actual content
    #     }
        
    #     # Step 4: Retriever is now used directly in build_chain_input to avoid parallel execution
        
    #     # Step 5: Build system prompt with form context (but NOT with RAG context yet)
    #     # RAG context will be dynamically added per query via {context} placeholder
    #     base_system_prompt = self._get_system_prompt(form_context)
        
    #     # Escape curly braces to prevent template variable errors
    #     escaped_system_prompt = base_system_prompt.replace("{", "{{").replace("}", "}}")
        
    #     # Build prompt template
    #     prompt = ChatPromptTemplate.from_messages([
    #         ("system", escaped_system_prompt + "\n\n### Knowledge Base Context:\n{context}" + uploaded_doc_content),
    #         MessagesPlaceholder(variable_name="chat_history"),
    #         ("human", "{user_query}")
    #     ])

    #     # Step 6: Build chain - restructured to avoid parallel execution that causes thread exhaustion
    #     def extract_query_content(x):
    #         """Extract user query from input - handles both string and message objects"""
    #         user_input = x.get("user_query", "")
    #         if hasattr(user_input, 'content'):
    #             return user_input.content
    #         return str(user_input)
        
    #     # def build_chain_input(x):
    #     #     """Build chain input sequentially to avoid parallel execution"""
    #     #     query = extract_query_content(x)
    #     #     # Retrieve context synchronously (not in parallel)
    #     #     context = self.format_context(retriever.invoke(query))
    #     #     chat_history = x.get("chat_history", [])
    #     #     return {
    #     #         "context": context,
    #     #         "user_query": query,
    #     #         "chat_history": chat_history
    #     #     }
    #     def build_chain_input(x):
    #         query = extract_query_content(x)
    #         docs = retriever.invoke(query)  # Sequential, not parallel
    #         return {
    #             "context": self.format_context(docs),
    #             "user_query": query,
    #             "chat_history": x.get("chat_history", [])
    #         }
        
        
    #     chain = (
    #         RunnableLambda(build_chain_input)
    #         | prompt 
    #         | self.llm
    #     )

    #     # Step 7: Wrap with message history
    #     conversational_chain = RunnableWithMessageHistory(
    #         chain,
    #         self.get_session_history,
    #         input_messages_key="user_query",
    #         history_messages_key="chat_history"
    #     )
        

    #     # Step 8: For initial message, retrieve document overview to inform the greeting
    #     # Check if this is the first message (no chat history)
    #     chat_history = self.get_session_history(session_id)
    #     is_first_message = len(chat_history.messages) == 0 if hasattr(chat_history, 'messages') else True
        
    #     # If first message and document is uploaded, enhance the query with document overview
    #     if is_first_message and document_uploaded:
    #         # Retrieve document overview to understand what it's about
    #         try:
    #             overview_query = "What is this document about? Provide a brief summary of the main topics, themes, and content."
    #             overview_docs = retriever.invoke(overview_query)
    #             if overview_docs:
    #                 # Extract key topics from the document
    #                 doc_summary = "\n".join([doc.page_content[:200] for doc in overview_docs[:3]])  # First 200 chars of top 3 chunks
    #                 # Enhance user query to include document context for better initial response
    #                 enhanced_query = f"{user_query}\n\n[Document Context: The uploaded document covers: {doc_summary}...]"
    #             else:
    #                 enhanced_query = user_query
    #         except Exception as e:
    #             teacher_logger.warning(f"Could not retrieve document overview: {str(e)}")
    #             enhanced_query = user_query
    #     else:
    #         enhanced_query = user_query
        
    #     # Step 8: Invoke with configuration to prevent thread exhaustion
    #     # Use simple config to avoid parallel execution issues

    #     run_config = RunnableConfig(
    #         configurable={"session_id": session_id},
    #         max_concurrency=1,
    #         recursion_limit=3
    #     )
        
    #     response_message = conversational_chain.invoke(
    #         {"user_query": enhanced_query},
    #         config=run_config
    #     )
        # run_config = RunnableConfig(
        #     configurable={"session_id": session_id}
        # )
        # response_message = conversational_chain.invoke(
        #     {"user_query": enhanced_query},
        #     config=run_config
        # )
        
        response_text = response_message.content if hasattr(response_message, 'content') else str(response_message)
        
        # Step 9: Check if complete lesson has been generated
        try:
            lesson_check = check_lesson_response(response_text, self.api_key)
            complete_lesson_status = lesson_check.complete_lesson
            teacher_logger.info(f"Lesson completion check: {complete_lesson_status}")
        except Exception as e:
            teacher_logger.warning(f"Error checking lesson completion status: {str(e)}")
            complete_lesson_status = "no"
        
        return InteractiveChatResponse(
            ai_response=response_text,
            complete_lesson=complete_lesson_status
        )
    # def interactive_chat(
    #     self, 
    #     lesson_id: int, 
    #     user_query: str, 
    #     session_id: str = None,
    #     subject: str = None,
    #     grade_level: str = None,
    #     focus_area: str = None,
    #     document_uploaded: bool = False,
    #     document_filename: str = None
    # ) -> InteractiveChatResponse:
    #     """Interactive chat with Prof. Potter for lesson creation"""
        
    #     # Use lesson_id as session_id to maintain context per lesson
    #     if not session_id:
    #         session_id = f"lesson_{lesson_id}"
        
    #     # Store form data in session context for use in prompts
    #     form_context = {
    #         'subject': subject or focus_area,  # Use focus_area as subject if subject not provided
    #         'grade_level': grade_level,
    #         'document_uploaded': document_uploaded,
    #         'document_filename': document_filename
    #     }
        
    #     # Step 1: Load vector DB
    #     vector_db = FAISS.load_local(
    #         "vector_store.faiss",
    #         self.rag_service.embeddings,
    #         allow_dangerous_deserialization=True
    #     )
    #     retriever = vector_db.as_retriever(search_type="similarity", search_kwargs={"k": 5})
        
    #     # Step 2: Create retriever runnable that extracts user_query from dict
    #     # Handle both string and message object inputs
    #     def get_query_for_retrieval(x):
    #         """Extract user query for retrieval - handles both string and message objects"""
    #         user_input = x.get("user_query", "")
    #         if hasattr(user_input, 'content'):
    #             return user_input.content
    #         return str(user_input)
        
    #     retriever_runnable = RunnableLambda(
    #         lambda x: self.format_context(
    #             retriever.invoke(get_query_for_retrieval(x))  # Extract query content for retrieval
    #         )
    #     )
        
    #     # Step 3: Build prompt (include context in the system prompt)
    #     # Get initial context for system prompt (returns list of Document objects)
    #     initial_context_docs = retriever.invoke(user_query)
    #     # Format context to string for system prompt (don't escape yet, will escape in system_prompt)
    #     initial_context_str = self._format_context_for_system_prompt(initial_context_docs)
    #     # Pass form context to system prompt
    #     system_prompt = self._get_system_prompt(initial_context_str, form_context)
    #     # Escape curly braces in system_prompt content to prevent template variable errors
    #     # This escapes any braces in the context content like {'producer'} that could be misinterpreted
    #     # The {context} variable is added after escaping, so it remains a valid template variable
    #     escaped_system_prompt = system_prompt.replace("{", "{{").replace("}", "}}")
    #     # Build prompt template with:
    #     # - System prompt with context
    #     # - MessagesPlaceholder for chat_history (automatically populated by RunnableWithMessageHistory)
    #     # - Current user query
    #     prompt = ChatPromptTemplate.from_messages([
    #         ("system", escaped_system_prompt + "\n\nContext from knowledge base:\n{context}"),
    #         MessagesPlaceholder(variable_name="chat_history"),  # Will contain all previous Q&A pairs
    #         ("human", "{user_query}")  # Current user query
    #     ])

    #     # Step 4: Build chain
    #     # RunnableWithMessageHistory will automatically:
    #     # 1. Convert user_query string to HumanMessage and add to history
    #     # 2. Inject full history into chat_history for MessagesPlaceholder
    #     # 3. Pass the original user_query string to the chain for template
    #     def extract_query_content(x):
    #         """Extract user query from input - handles both string and message objects"""
    #         user_input = x.get("user_query", "")
    #         if hasattr(user_input, 'content'):
    #             return user_input.content
    #         return str(user_input)
        
    #     chain = (
    #         {
    #             "context": retriever_runnable,  # Retrieves context based on user_query
    #             "user_query": extract_query_content,  # Extract query content (handles message or string)
    #             "chat_history": lambda x: x.get("chat_history", [])  # Pass chat_history (populated by RunnableWithMessageHistory)
    #         }
    #         | prompt 
    #         | self.llm
    #     )

    #     # Step 5: Wrap with message history
    #     # RunnableWithMessageHistory automatically:
    #     # - Converts user_query string to HumanMessage
    #     # - Adds it to session history
    #     # - Injects full history into chat_history key before passing to chain
    #     conversational_chain = RunnableWithMessageHistory(
    #         chain,
    #         self.get_session_history,
    #         input_messages_key="user_query",  # Key in input dict that contains the message
    #         history_messages_key="chat_history"  # Key in chain input where history will be injected
    #     )

    #     # Step 6: Invoke
    #     # Pass user_query as string - RunnableWithMessageHistory will handle conversion and history
    #     response_message = conversational_chain.invoke(
    #         {"user_query": user_query},  # Pass as string, will be converted to HumanMessage
    #         config={"configurable": {"session_id": session_id}}
    #     )
        
    #     response_text = response_message.content if hasattr(response_message, 'content') else str(response_message)
        
    #     # Step 7: Check if complete lesson has been generated
    #     # Use structured output to determine if the response indicates a complete lesson
    #     try:
    #         lesson_check = check_lesson_response(response_text, self.api_key)
            
    #         complete_lesson_status = lesson_check.complete_lesson
    #         teacher_logger.info(f"Lesson completion check: {complete_lesson_status}")
    #     except Exception as e:
    #         teacher_logger.warning(f"Error checking lesson completion status: {str(e)}")
    #         # Default to "no" if check fails (assume it's still in draft/outline stage)
    #         complete_lesson_status = "no"
        
    #     # Return structured response
    #     return InteractiveChatResponse(
    #         ai_response=response_text,
    #         complete_lesson=complete_lesson_status
    #     )
            
        # except Exception as e:
        #     teacher_logger.error(f"Error in conversational chain: {e}")
        #     # Return fallback response
        #     return InteractiveChatResponse(
        #         ai_response=f"I encountered an error. Let's start fresh. Could you please tell me what lesson you'd like to create?",
        #         state="start",
        #         should_generate_lesson=False
        #     )
    
    # def _parse_response(self, response_text: str, current_state: str) -> InteractiveChatResponse:
    #     """Parse LLM response and extract structured data"""
        
    #     # Try to extract JSON if present
    #     json_match = re.search(r'\{[^}]*"ai_response"[^}]*\}', response_text, re.DOTALL)
    #     if json_match:
    #         try:
    #             data = json.loads(json_match.group())
    #             return InteractiveChatResponse(**data)
    #         except:
    #             pass
        
    #     # Determine next state based on content and current state
    #     next_state = self._determine_next_state(response_text, current_state)
        
    #     # Check if lesson should be generated
    #     should_generate = (
    #         next_state == "completed" or
    #         "here is the complete lesson" in response_text.lower() or
    #         "here's the complete lesson" in response_text.lower()
    #     )
        
    #     return InteractiveChatResponse(
    #         ai_response=response_text.strip(),
    #         state=next_state,
    #         should_generate_lesson=should_generate
    #     )
    
    # def _determine_next_state(self, response_text: str, current_state: str) -> str:
    #     """Determine next state based on response content"""
    #     response_lower = response_text.lower()
        
    #     if current_state == "start":
    #         if "have you uploaded" in response_lower or "upload" in response_lower:
    #             return "gathering_requirements"
    #         return "gathering_requirements"
        
    #     elif current_state == "gathering_requirements":
    #         if "outline" in response_lower or "structure" in response_lower:
    #             return "awaiting_outline_approval"
    #         return "gathering_requirements"
        
    #     elif current_state == "awaiting_outline_approval":
    #         if "generate the complete lesson" in response_lower or "full lesson" in response_lower:
    #             return "ready_to_generate"
    #         return "refining_outline"
        
    #     elif current_state == "refining_outline":
    #         if "ready to generate" in response_lower or "generate" in response_lower:
    #             return "ready_to_generate"
    #         return "refining_outline"
        
    #     elif current_state == "ready_to_generate":
    #         return "completed"
        
    #     elif current_state == "completed":
    #         return "completed"
        
    #     return current_state
    

    def _get_system_prompt(self, rag_context: str = "", form_context: dict = None) -> str:
        """Generate a single unified system prompt for Prof. Potter."""
        
        if form_context is None:
            form_context = {}
        
        # Note: rag_context is typically empty here and filled by {context} placeholder
        # Only include it if explicitly provided (rare case)
        context_section = f"\n\n### RELEVANT CONTEXT FROM KNOWLEDGE BASE:\n{rag_context}\n" if rag_context else ""
        
        # Build form context section
        form_context_section = ""
        # Only show subject if it's not "Other" or empty
        subject = form_context.get('subject', '')
        if subject and subject.lower() not in ['other', 'none', '']:
            form_context_section += f"\n**Subject/Topic:** {subject}\n"
        if form_context.get('grade_level'):
            form_context_section += f"**Grade Level:** {form_context['grade_level']}\n"
        # Document is already uploaded - don't ask about it, just note it's available
        if form_context.get('document_uploaded'):
            form_context_section += f"**Document:** {form_context.get('document_filename', 'file')} (already uploaded and processed)\n"
            form_context_section += "**CRITICAL INSTRUCTION:** The document has been uploaded and processed. You have access to its full content through the knowledge base context. Use this content to understand what the document is about and help create lessons from it.\n"
            form_context_section += "**WHEN FACULTY REQUESTS LESSON GENERATION:** DO NOT ask about uploading documents or textbooks. The document is already available. Proceed immediately to create the lesson using the document content from the knowledge base context. Acknowledge: 'Perfect! I'll create a lesson plan using the document you've already uploaded ([filename]). Let me analyze the content...'\n"
        
        if form_context_section:
            form_context_section = f"\n\n📋 LESSON FORM INFORMATION:\n{form_context_section}\n"

        unified_prompt = f"""
              ROLE:
You are Prof. Potter, an expert education assistant helping Faculty/Teachers create lessons from their already uploaded document, which may include text, images, tables, charts, or graphs.
The document has been uploaded before this chat started — never ask for uploads.

INTERACTION MODES
MODE 1: General Conversation

Warm, natural tone (≤200 words)

Answer using the document when possible (including text, images & tables)

If the question is NOT in the document:
“I cannot find [topic] in your document. Your document covers [actual topics].”

MODE 2: Lesson Planning

Create structured lessons using ONLY the document content

Step-by-step, ≤150 words per step

Faculty must choose a topic first

Respect all prerequisites & confirmation steps

CRITICAL RULES
RULE 1: Document Is Already Uploaded

NEVER ask for the document

ALWAYS check the RELEVANT CONTEXT FROM KNOWLEDGE BASE

Use found info.

Only say “I cannot find…” if the context truly lacks it.

RULE 2: When Faculty Says “Generate Lesson”

Respond:
“Perfect! Your document covers: [topics]. Which topic should I focus on?”
→ Then WAIT.

RULE 3: Wait After Every Question

Do not continue until Faculty responds.

RULE 4: Check Context First — INCLUDING Images, Tables & Graphs

Tables must be read and used in answers

Charts/graphs must be interpreted (trends, values, axes, message)

Images must be described and analyzed

If relevant info appears anywhere (text or visual), you MUST use it

MULTIMODAL CONTENT RULES
1. For Images

Describe them clearly using structured analysis:

This image shows…

All visible elements

Actions, labels, diagrams

Educational meaning or purpose

If the user asks: extract info from the image

2. For Graphs/Charts

Explain:

Type (bar, pie, line, etc.)

What the chart is telling

Values, labels, axes

Trends, comparisons

Key insight

3. For Tables

Use them as factual data:

Read rows/columns

Provide exact numbers if asked

Summarize patterns

Cite table info when answering questions

FIRST MESSAGE (Greeting)

“Hello! I’ve reviewed your document about [topics from context including visuals, tables & charts]. Would you like me to create a lesson plan from this content?”

LESSON GENERATION FLOW (Improved)

Faculty: “Generate lesson”
You: “Perfect! Your document covers: [topics]. Which topic should I focus on?”

Faculty chooses topic
You: “For [topic], students need [prerequisites]. Include them?”

Faculty confirms
You: Generate lesson step-by-step (≤150 words per step)

After lesson:
“Does this work for your students?”

HANDLING MULTIPLE QUESTIONS

User may ask several things at once

Answer each according to context

Remind them gently:
“If you’d like, I can also generate a lesson from the document.”

WHEN INFORMATION IS NOT FOUND

Only after checking carefully:
“I cannot find [topic] in your document. Your document covers [topics]. Would you like a lesson from the available content or an answer from general knowledge?”

EQUATION TEACHING RULES

Explain each term

Explain operations

Build full equation

Add simple real-world example

Ask if they want to proceed

QUALITY CHECKLIST

Before responding, ensure:

✔ Checked the RELEVANT CONTEXT FROM KNOWLEDGE BASE

✔ Included images, charts, tables when relevant

✔ No fabricated info

✔ Warm, clear tone

✔ Word limits respected

✔ Stop and wait for user when required

NEVER DO THIS

❌ Ask for document
❌ Ignore images/tables/charts
❌ Invent missing content
❌ Continue lesson without confirmation
❌ Exceed word limits

ALWAYS DO THIS

✔ Use all context (text + visuals + data)
✔ Interpret tables & graphs correctly
✔ Offer lesson generation proactively but only proceed when user chooses
✔ Stay warm, supportive, professional
✔ Help faculty stay aware of lesson planning options

{form_context_section}{rag_context}

CRITICAL REMINDER:
The "RELEVANT CONTEXT FROM KNOWLEDGE BASE" section above contains information retrieved from the document based on the user's query.
- If the context contains information related to the question → ANSWER using that information
- Do NOT say "I cannot find" if the information is actually in the provided context
- Only say "I cannot find" if the context truly does NOT contain relevant information
- The context is already filtered and relevant to the query, so USE IT!

TONE: Warm, professional, encouraging, clear, natural.
REMEMBER: Document uploaded BEFORE chat. The context provided above is from the document - USE IT if it contains relevant information!"""
        return unified_prompt.strip()

    def _llm_responce(self, rag_prompt: str, lesson_details: Optional[Dict[str, str]] = None) -> str:
        """
        Generate lesson response using RAG prompt (for large documents)
        
        Args:
            rag_prompt: Pre-formatted prompt with retrieved context
            lesson_details: Additional lesson details
            
        Returns:
            String response from LLM
        """
        teacher_logger.info("=== AI LESSON GENERATION WITH RAG STARTED ===")
        teacher_logger.info(f"RAG prompt length: {len(rag_prompt)} characters")
        
        try:
            # Use the RAG prompt directly with the LLM
            response = self.llm.invoke(rag_prompt)
            teacher_logger.info("LLM response received")
            
            # Extract content from AIMessage
            if hasattr(response, 'content'):
                response_content = response.content
            else:
                response_content = str(response)
            
            teacher_logger.info(f"Response content length: {len(response_content)} characters")
            teacher_logger.info("=== AI LESSON GENERATION WITH RAG COMPLETED ===")
            
            return response_content
                
        except Exception as e:
            teacher_logger.error(f"Error in RAG lesson generation: {str(e)}")
            return f"Error generating lesson: {str(e)}"


        

    def _generate_structured_lesson(self, text: str, lesson_details: Optional[Dict[str, str]] = None) -> str:
        """Generate lesson response based on user intent."""
        teacher_logger.info("=== AI LESSON GENERATION STARTED ===")
        teacher_logger.info(f"Text length: {len(text)} characters")
        
        try:
            max_chars = 1000000
            if len(text) > max_chars:
                text = text[:max_chars] + "..."
                teacher_logger.info(f"Text truncated to {max_chars} characters")
                logger.info(f"Text truncated to {max_chars} characters")

            user_prompt = ""
            grade_level = ""
            focus_area = ""

            if lesson_details:
                user_prompt = lesson_details.get("lesson_prompt", "").strip()
                grade_level = lesson_details.get("grade_level", "")
                focus_area = lesson_details.get("focus_area", "")

            teacher_logger.info(f"User prompt: {user_prompt}")
            teacher_logger.info(f"Grade level: {grade_level}")
            teacher_logger.info(f"Focus area: {focus_area}")

            return self._generate_direct_answer(text, user_prompt)

        except Exception as e:
            teacher_logger.error(f"Structured lesson generation failed: {str(e)}")
            logger.error(f"Error generating structured lesson: {str(e)}", exc_info=True)
            return f"Error generating lesson: {str(e)}"

    # def _user_wants_lesson_plan(self, user_prompt: str) -> bool:
    #     """Determine if user wants a lesson plan or just an answer"""
    #     if not user_prompt:
    #         return True  # Default to lesson plan if no prompt
        
    #     # Keywords that indicate user wants a lesson plan
    #     lesson_keywords = [
    #         "generate lesson", "create lesson", "make lesson", "lesson plan",
    #         "teaching plan", "educational plan", "curriculum", "syllabus",
    #         "lesson from", "lesson based on", "teach this", "explain as lesson"
    #     ]
        
    #     # Keywords that indicate user wants a direct answer
    #     question_keywords = [
    #         "what is", "how does", "explain", "tell me", "describe",
    #         "why", "when", "where", "which", "who", "how",
    #         "question", "answer", "help me understand"
    #     ]
        
    #     prompt_lower = user_prompt.lower()
        
    #     # Check for lesson keywords first
    #     for keyword in lesson_keywords:
    #         if keyword in prompt_lower:
    #             return True
        
    #     # Check for question keywords
    #     for keyword in question_keywords:
    #         if keyword in prompt_lower:
    #             return False
        
    #     # Default to lesson plan if unclear
    #     return True

    def _generate_lesson_plan(self, text: str, user_prompt: str, grade_level: str, focus_area: str) -> Dict[str, Any]:
        """Generate a comprehensive lesson plan"""
        teacher_logger.info("=== LESSON PLAN GENERATION STARTED ===")
        teacher_logger.info(f"Text length: {len(text)}")
        teacher_logger.info(f"User prompt: {user_prompt}")
        teacher_logger.info(f"Grade level: {grade_level}")
        teacher_logger.info(f"Focus area: {focus_area}")
        
        try:
            # Create Pydantic output parser
            teacher_logger.info("Setting up Pydantic output parser")
            parser = PydanticOutputParser(pydantic_object=LessonResponse)

            # Create prompt template for lesson plan
            teacher_logger.info("Creating prompt template for lesson plan")
            prompt_template = PromptTemplate(
                template="""You are an expert teacher. Create a comprehensive lesson plan based on the document content and user request.

Document Content:
{text}

User Request:
{user_prompt}

Grade Level: {grade_level}
Focus Area: {focus_area}

Create a complete lesson plan with:
- Clear title and summary
- Learning objectives
- Background prerequisites
- Structured sections with detailed content
- Creative activities for students
- Assessment questions
- Teacher notes

Use language appropriate for grade level: {grade_level}
Focus on subject area: {focus_area}

{format_instructions}""",
                input_variables=["text", "user_prompt", "grade_level", "focus_area"],
                partial_variables={"format_instructions": parser.get_format_instructions()}
            )

            # Create chain with LLM and parser
            teacher_logger.info("Creating LLM chain")
            chain = prompt_template | self.llm | parser

            # Invoke the chain
            teacher_logger.info("Invoking LLM for lesson plan generation")
            result = chain.invoke({
                "text": text,
                "user_prompt": user_prompt,
                "grade_level": grade_level,
                "focus_area": focus_area
            })
            
            teacher_logger.info(f"LLM response received: {result.response_type}")
            logger.info(f"Generated lesson plan: {result.response_type}")
            
            if result.response_type == "lesson_plan" and result.answer:
                lesson_dict = result.answer.dict()
                teacher_logger.info(f"Lesson plan generated successfully with {len(lesson_dict.get('sections', []))} sections")
                teacher_logger.info(f"Learning objectives: {len(lesson_dict.get('learning_objectives', []))}")
                teacher_logger.info(f"Creative activities: {len(lesson_dict.get('creative_activities', []))}")
                teacher_logger.info(f"Quiz questions: {len(lesson_dict.get('assessment_quiz', []))}")
                teacher_logger.info("=== LESSON PLAN GENERATION COMPLETED ===")
                
                return {
                    "response_type": "lesson_plan",
                    "lesson": lesson_dict
                }
            else:
                teacher_logger.warning("LLM didn't generate lesson plan, using fallback")
                logger.warning("LLM didn't generate lesson plan, using fallback")
                return self._create_fallback_lesson(text)
                
        except Exception as e:
            teacher_logger.error(f"Lesson plan generation failed: {str(e)}")
            logger.error(f"Error generating lesson plan: {str(e)}")
            return self._create_fallback_lesson(text)

    def _generate_direct_answer(self, text: str, user_prompt: str) -> str:
        """Generate a direct detailed answer to user's question"""
        teacher_logger.info("=== DIRECT ANSWER GENERATION STARTED ===")
        teacher_logger.info(f"Text length: {len(text)}")
        teacher_logger.info(f"User question: {user_prompt}")
        
        try:
            # Create a simple prompt for direct answers
            teacher_logger.info("Creating prompt for direct answer")
            answer_prompt = f"""You are a helpful teacher. Answer the user's question based on the document content.

Document Content:
{text}

User Question:
{user_prompt}

Provide a detailed, comprehensive answer that directly addresses the user's question. Use information from the document content to support your answer. Be thorough and educational.

Answer:"""

            teacher_logger.info("Invoking LLM for direct answer")
            response = self.llm.invoke(answer_prompt)
            
            # Extract content from AIMessage
            if hasattr(response, 'content'):
                answer_text = response.content.strip()
            else:
                answer_text = str(response).strip()
            
            teacher_logger.info(f"Direct answer generated successfully - length: {len(answer_text)}")
            teacher_logger.info("=== DIRECT ANSWER GENERATION COMPLETED ===")
            logger.info(f"Generated direct answer (length: {len(answer_text)})")
            
            return answer_text
            
        except Exception as e:
            teacher_logger.error(f"Direct answer generation failed: {str(e)}")
            logger.error(f"Error generating direct answer: {str(e)}")
            return f"I apologize, but I encountered an error while processing your question: {str(e)}"

    def _create_fallback_lesson(self, text: str) -> Dict[str, Any]:
        """Create a fallback lesson when AI generation fails"""
        return {
            "response_type": "lesson_plan",
            "lesson": {
                "title": "Generated Lesson",
                "summary": "A lesson generated from your content.",
                "learning_objectives": ["Understand the key concepts"],
                "key_concepts": ["Key concepts from the content"],
                "background_prerequisites": ["Basic understanding of the topic"],
                "sections": [
                    {
                        "heading": "Introduction",
                        "content": text[:1000] + "..." if len(text) > 1000 else text
                    }
                ],
                "creative_activities": [],
                "stem_equations": [],
                "assessment_quiz": [],
                "teacher_notes": ["Review the content before teaching"]
            }
        }

    def _extract_sections_from_prompt(self, text: str, prompt: str) -> str:
        """Extract specific sections based on user prompt"""
        # Simple keyword-based extraction
        # This could be enhanced with more sophisticated NLP
        keywords = prompt.lower().split()
        sentences = text.split('.')
        relevant_sentences = []
        
        for sentence in sentences:
            if any(keyword in sentence.lower() for keyword in keywords):
                relevant_sentences.append(sentence.strip())
        
        return '. '.join(relevant_sentences) if relevant_sentences else ""

    def _extract_lesson_text_for_rag(self, lesson_data: Dict[str, Any]) -> str:
        """Extract all text content from lesson data for RAG storage"""
        try:
            text_parts = []
            
            # Add title and summary
            if lesson_data.get('title'):
                text_parts.append(f"Title: {lesson_data['title']}")
            if lesson_data.get('summary'):
                text_parts.append(f"Summary: {lesson_data['summary']}")
            
            # Add learning objectives
            if lesson_data.get('learning_objectives'):
                text_parts.append("Learning Objectives:")
                for obj in lesson_data['learning_objectives']:
                    text_parts.append(f"- {obj}")
            
            # Add sections
            if lesson_data.get('sections'):
                for section in lesson_data['sections']:
                    if section.get('heading'):
                        text_parts.append(f"\n{section['heading']}:")
                    if section.get('content'):
                        text_parts.append(section['content'])
            
            # Add key concepts
            if lesson_data.get('key_concepts'):
                text_parts.append("\nKey Concepts:")
                for concept in lesson_data['key_concepts']:
                    text_parts.append(f"- {concept}")
            
            # Add background prerequisites
            if lesson_data.get('background_prerequisites'):
                text_parts.append("\nBackground Prerequisites:")
                for prereq in lesson_data['background_prerequisites']:
                    text_parts.append(f"- {prereq}")
            
            # Add creative activities
            if lesson_data.get('creative_activities'):
                text_parts.append("\nCreative Activities:")
                for activity in lesson_data['creative_activities']:
                    if activity.get('name'):
                        text_parts.append(f"\n{activity['name']}:")
                    if activity.get('description'):
                        text_parts.append(activity['description'])
                    if activity.get('learning_purpose'):
                        text_parts.append(f"Purpose: {activity['learning_purpose']}")
            
            # Add STEM equations
            if lesson_data.get('stem_equations'):
                text_parts.append("\nSTEM Equations:")
                for eq in lesson_data['stem_equations']:
                    if eq.get('equation'):
                        text_parts.append(f"Equation: {eq['equation']}")
                    if eq.get('complete_equation_significance'):
                        text_parts.append(f"Significance: {eq['complete_equation_significance']}")
            
            # Add assessment quiz
            if lesson_data.get('assessment_quiz'):
                text_parts.append("\nAssessment Quiz:")
                for q in lesson_data['assessment_quiz']:
                    if q.get('question'):
                        text_parts.append(f"Q: {q['question']}")
                    if q.get('answer'):
                        text_parts.append(f"A: {q['answer']}")
            
            # Add teacher notes
            if lesson_data.get('teacher_notes'):
                text_parts.append("\nTeacher Notes:")
                for note in lesson_data['teacher_notes']:
                    text_parts.append(f"- {note}")
            
            return "\n".join(text_parts)
            
        except Exception as e:
            teacher_logger.error(f"Error extracting lesson text for RAG: {str(e)}")
            return ""

    def _store_original_document_rag(self, filename: str):
        """Store the original document's RAG service for AI review"""
        try:
            # Create a unique key for this document
            document_key = f"document_{filename}_{hash(filename) % 10000}"
            
            # Store the current RAG service instance (which contains the original document)
            if hasattr(self.rag_service, 'vector_store') and self.rag_service.vector_store:
                # Create a copy of the RAG service to avoid overwriting
                from copy import deepcopy
                rag_service_copy = deepcopy(self.rag_service)
                
                self.lesson_vector_stores[document_key] = {
                    'rag_service': rag_service_copy,
                    'filename': filename,
                    'content': 'original_document',  # Mark as original document
                    'type': 'original_document'
                }
                teacher_logger.info(f"Original document RAG service stored with key: {document_key}")
                teacher_logger.info(f"Vector store contains {len(rag_service_copy.documents)} chunks")
            else:
                teacher_logger.warning("No vector store available to store")
                
        except Exception as e:
            teacher_logger.error(f"Error storing original document RAG: {str(e)}")

    def _store_lesson_in_vector_db(self, lesson_content: str, filename: str):
        """Store lesson content in vector database for AI review"""
        try:
            # Create a unique key for this lesson
            lesson_key = f"lesson_{filename}_{hash(lesson_content) % 10000}"
            
            # Create documents from lesson content
            from langchain_core.documents import Document
            documents = [Document(page_content=lesson_content, metadata={"lesson_key": lesson_key, "filename": filename})]
            
            # Process with RAG service
            rag_result = self.rag_service.process_document(documents, filename)
            if 'error' not in rag_result:
                # Store the RAG service instance for this lesson
                self.lesson_vector_stores[lesson_key] = {
                    'rag_service': self.rag_service,
                    'filename': filename,
                    'content': lesson_content
                }
                teacher_logger.info(f"Lesson stored in vector DB with key: {lesson_key}")
            else:
                teacher_logger.error(f"Failed to store lesson in vector DB: {rag_result['error']}")
                
        except Exception as e:
            teacher_logger.error(f"Error storing lesson in vector DB: {str(e)}")

    def review_lesson_with_rag(self, lesson_content: str, user_prompt: str, filename: str = "") -> str:
        """Review lesson content using RAG to retrieve relevant information"""
        try:
            teacher_logger.info("=== RAG-BASED LESSON REVIEW STARTED ===")
            teacher_logger.info(f"User prompt: {user_prompt}")
            teacher_logger.info(f"Filename: {filename}")
            
            # Try to find existing vector store for this lesson
            lesson_key = None
            for key, store_data in self.lesson_vector_stores.items():
                if store_data['filename'] == filename or store_data['content'] == lesson_content:
                    lesson_key = key
                    break
            
            if lesson_key and lesson_key in self.lesson_vector_stores:
                teacher_logger.info(f"Using existing vector store: {lesson_key}")
                rag_service = self.lesson_vector_stores[lesson_key]['rag_service']
            else:
                teacher_logger.info("Creating new vector store for lesson review")
                # Create new RAG service for this lesson
                from langchain_core.documents import Document
                documents = [Document(page_content=lesson_content, metadata={"filename": filename})]
                rag_service = RAGService()
                rag_result = rag_service.process_document(documents, filename)
                
                if 'error' in rag_result:
                    teacher_logger.error(f"Failed to create vector store: {rag_result['error']}")
                    # Fallback to regular improvement
                    return self.improve_lesson_content(0, lesson_content, user_prompt)
            
            # Retrieve relevant chunks
            relevant_chunks = rag_service.retrieve_relevant_chunks(user_prompt, k=5)
            if not relevant_chunks:
                teacher_logger.warning("No relevant chunks found, using full content")
                relevant_chunks = [Document(page_content=lesson_content, metadata={})]
            
            # Create RAG prompt
            rag_prompt = rag_service.create_rag_prompt(user_prompt, relevant_chunks)
            
            # Generate response using RAG
            teacher_logger.info("Generating RAG-based response")
            response = self.llm.invoke(rag_prompt)

            response=response.content
            
            # response_content = response.content if hasattr(response, 'content') else str(response)

            
            
            # Check if the response is in JSON format and extract the actual content
            # try:
            #     import json
            #     parsed_response = json.loads(response_content)
            #     if isinstance(parsed_response, dict) and 'answer' in parsed_response:
            #         response_content = parsed_response['answer']
            #     elif isinstance(parsed_response, dict) and 'response' in parsed_response:
            #         response_content = parsed_response['response']
            #     elif isinstance(parsed_response, dict) and 'content' in parsed_response:
            #         response_content = parsed_response['content']
            # except (json.JSONDecodeError, AttributeError):
            #     pass
            
            teacher_logger.info("=== RAG-BASED LESSON REVIEW COMPLETED ===")
            return response
            
        except Exception as e:
            teacher_logger.error(f"Error in RAG-based lesson review: {str(e)}")
            # Fallback to regular improvement
            return "no relevant info found"
            # return self.improve_lesson_content(0, lesson_content, user_prompt)

    def improve_lesson_content(self, lesson_id: int, current_content: str, improvement_prompt: str = "") -> str:
        """Improve lesson content using AI based on user prompt"""
        try:
            # Create improvement prompt
            #fetch the original content from lesson id
            original_content = LessonModel.get_lesson_by_id(lesson_id)['original_content']
            if improvement_prompt:
                
                prompt = f"""You are an expert assistant. Improve or modify the following content based on the user's specific request or if user ask any question from the content also respond the answer from the content as well.

            User's Request:
            {improvement_prompt}

            Original Content:
            {original_content}

            Current Content:
            {current_content}

            CRITICAL INSTRUCTIONS:

            1. CONTENT SOURCES:
            - You have access to BOTH the original and current versions of the content
            - Use whichever version (or combination of both) best fulfills the user's request
            - If user wants to add back something from original: Pull from original content
            - If user wants to keep recent changes: Use current content
            - If user wants the best of both: Intelligently combine information from both versions
            - Default: Use current content as the primary base, but reference original if it helps

            2. MANDATORY FORMAT - Always use proper markdown formatting:
            - Use ## for main headers, ### for subheaders
            - Use **bold** for key terms and important concepts
            - Use - or * for bullet points
            - Use 1. 2. 3. for numbered lists
            - Use > for quotes or callouts
            - Use `code` for technical terms when appropriate
            - EXCEPTION: If user requests "one line" or "single line", output everything as ONE CONTINUOUS SENTENCE with NO line breaks, but still use **bold** and *italic* for emphasis

            3. CONTENT MODIFICATIONS based on user request:
            - "one line" or "single line" or "give me response in one line": 
                * MUST condense EVERYTHING into literally ONE SENTENCE
                * NO line breaks, NO bullet points, NO numbered lists, NO headers
                * Use **bold** for key terms within that single sentence
                * Example: "**John Doe** is a software engineer specializing in web development, cloud computing, and AI integration with 5 years of experience building scalable applications."
            
            - "concise" or "summary": 2-4 sentences with markdown formatting
            
            - "detailed" or "expanded": Multiple paragraphs with full markdown structure (headers, bullets, bold terms)
            
            - "structured": Organize with clear headers, subheaders, and bullet points
            
            - "combine both versions": Merge the best elements from original and current
            
            - "use original" or "revert to original": Base response primarily on original content
            
            - Any other request: Apply it flexibly using information from both versions as needed

            4. ABSOLUTE RULES:
            - If user says "one line" → Output must be EXACTLY one continuous sentence, no exceptions
            - For all other requests: ALWAYS use proper markdown formatting
            - User's request is the TOP priority - use original, current, or both as needed to satisfy it
            - Return ONLY the improved content
            - NO meta-commentary like "Since the user requested..." or "I have condensed..."
            - NO explanations about what you did or why
            - Just return the final improved content directly

            ##always follow this instruction##
              when user ask the question from the content also give the answer consiley if user exlicity request for
              **one** line answer also give the answer in one line
              if user say gave me two line answer gove the answer in two line
              if the user say give me detailed answer give the answer in detailed manner
            """
            # Generate improved content using LLM
            response = self.llm.invoke(prompt)
            improved_content = response.content.strip()
            
            # Check if the response is in JSON format and extract the actual content
            try:
                import json
                # Try to parse as JSON
                parsed_response = json.loads(improved_content)
                # If it has the expected JSON structure (like chat responses), extract the answer
                if isinstance(parsed_response, dict) and 'answer' in parsed_response:
                    improved_content = parsed_response['answer']
                elif isinstance(parsed_response, dict) and 'response' in parsed_response:
                    improved_content = parsed_response['response']
                elif isinstance(parsed_response, dict) and 'content' in parsed_response:
                    improved_content = parsed_response['content']
            except (json.JSONDecodeError, AttributeError):
                # Not a JSON response, use as-is
                pass
            
            logger.info(f"Successfully improved lesson {lesson_id} content")
            return improved_content
            
        except Exception as e:
            logger.error(f"Error improving lesson content: {str(e)}")
            # Return original content if improvement fails
            return current_content

    def edit_lesson_with_prompt(self, lesson_text: str, user_prompt: str, filename: str = "") -> str:
        """Use RAG system for semantic chunk retrieval and editing"""
        try:
            teacher_logger.info("=== AI REVIEW WITH RAG STARTED ===")
            teacher_logger.info(f"User prompt: {user_prompt}")
            teacher_logger.info(f"Filename: {filename}")
            teacher_logger.info(f"Lesson text length: {len(lesson_text)}")
            
            # Try to find existing vector store for this lesson/document
            lesson_key = None
            for key, store_data in self.lesson_vector_stores.items():
                if (store_data['filename'] == filename or filename in store_data['filename']) and store_data.get('type') == 'original_document':
                    lesson_key = key
                    break

            #retced the relevnt chunk from saved faiss
            from langchain_community.vectorstores import FAISS
            from langchain_community.embeddings import HuggingFaceEmbeddings
            import tempfile
            import os
            # embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

            # vector_db = FAISS.load_local("vector_store.faiss", embeddings, allow_dangerous_deserialization=True)
            # relevant_docs = vector_db.similarity_search(user_prompt, k=10)
            # # relevant_texts = [doc.page_content for doc in relevant_docs]
            # # lesson_text=["\n the previous version context".join(lesson_text)]

            # lesson_doc = Document(page_content="\n the previous version context".join(lesson_text))
            # relevant_docs.append(lesson_doc)
            from langchain_core.documents import Document
            from langchain_community.vectorstores import FAISS
            from langchain_community.embeddings import HuggingFaceEmbeddings

            # Environment variables TQDM_DISABLE and TOKENIZERS_PARALLELISM are set at app startup
            embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={
                    'device': 'cpu',
                    'trust_remote_code': False
                },
                encode_kwargs={
                    'normalize_embeddings': False
                }
            )
            vector_db = FAISS.load_local("vector_store.faiss", embeddings, allow_dangerous_deserialization=True)

            # Retrieve relevant docs (these are already Document objects)
            relevant_docs = vector_db.similarity_search(user_prompt, k=10)

            # Add the full lesson text as another Document
            lesson_doc = Document(page_content="\n the previous version context".join(lesson_text))
            relevant_docs.append(lesson_doc)
            #combined with relevant chunks
            # relevant_texts=relevant_texts+lesson_text
            rag_service=RAGService()

            

            #combine
            rag_prompt = rag_service.create_rag_prompt(user_prompt, relevant_docs)
                
                # Generate response using RAG
            teacher_logger.info("Generating RAG-based response from original document")
            response = self.llm.invoke(rag_prompt)
            response_content = response.content
            return response_content
            
            if lesson_key and lesson_key in self.lesson_vector_stores:
                teacher_logger.info(f"Using existing original document vector store: {lesson_key}")
                rag_service = self.lesson_vector_stores[lesson_key]['rag_service']
                
                # Retrieve relevant chunks from original document
                relevant_chunks = rag_service.retrieve_relevant_chunks(user_prompt, k=10)
                if not relevant_chunks:
                    teacher_logger.warning("No relevant chunks found, using first few chunks")
                    relevant_chunks = rag_service.documents[:9]
                
                # Create RAG prompt with retrieved context
                rag_prompt = rag_service.create_rag_prompt(user_prompt, relevant_chunks)
                
                # Generate response using RAG
                teacher_logger.info("Generating RAG-based response from original document")
                response = self.llm.invoke(rag_prompt)
                
                if hasattr(response, 'content'):
                    response_content = response.content
                else:
                    response_content = str(response)
                
                teacher_logger.info("=== AI REVIEW WITH RAG COMPLETED ===")
                return response_content
                
            else:
                teacher_logger.info("No existing original document vector store found, creating new one from lesson text")

                print("-----------------------fallback responce---------------------")
                print(lesson_text)
                print("-----------------------fallback responce---------------------")
                # Fallback to creating vector store from lesson text
                return self._edit_lesson_with_fallback_rag(lesson_text, user_prompt)
                
        except Exception as e:
            teacher_logger.error(f"Error in RAG-based lesson editing: {str(e)}")
            # Fallback to simple editing
            return self._edit_lesson_simple(lesson_text, user_prompt)
    
    def _edit_lesson_with_fallback_rag(self, lesson_text: str, user_prompt: str) -> str:
        """Fallback method to create RAG from lesson text"""
        try:
            from langchain_community.vectorstores import FAISS
            from langchain_community.embeddings import HuggingFaceEmbeddings
            import tempfile
            import os

            # 1. Chunk the lesson (by paragraph)
            chunks = [p.strip() for p in lesson_text.split('\n\n') if p.strip()]
            if not chunks:
                return lesson_text

            # 2. Embed and store in FAISS
            # Environment variables TQDM_DISABLE and TOKENIZERS_PARALLELISM are set at app startup
            embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={
                    'device': 'cpu',
                    'trust_remote_code': False
                },
                encode_kwargs={
                    'normalize_embeddings': False
                }
            )
            
            with tempfile.TemporaryDirectory() as tmpdir:
                faiss_path = os.path.join(tmpdir, "faiss_index")
                vector_db = FAISS.from_texts(chunks, embeddings)
                vector_db.save_local(faiss_path)

                # 3. Retrieve relevant chunks for the user prompt
                vector_db = FAISS.load_local(faiss_path, embeddings, allow_dangerous_deserialization=True)
                relevant_docs = vector_db.similarity_search(user_prompt, k=min(5, len(chunks)))
                relevant_texts = [doc.page_content for doc in relevant_docs]

                # 4. Edit relevant chunks with LLM
                edited_chunks = {}
                for text in relevant_texts:
                    edit_prompt = (
                        f"You are an expert teacher assistant. Here is a lesson chunk:\n\n"
                        f"{text}\n\n"
                        f"User request: {user_prompt}\n\n"
                        f"Please return the revised lesson chunk."
                    )
                    response = self.llm.invoke(edit_prompt)
                    edited = response.content if hasattr(response, 'content') else str(response)
                    
                    # Check if the response is in JSON format and extract the actual content
                    try:
                        import json
                        parsed_response = json.loads(edited)
                        if isinstance(parsed_response, dict) and 'answer' in parsed_response:
                            edited = parsed_response['answer']
                        elif isinstance(parsed_response, dict) and 'response' in parsed_response:
                            edited = parsed_response['response']
                        elif isinstance(parsed_response, dict) and 'content' in parsed_response:
                            edited = parsed_response['content']
                    except (json.JSONDecodeError, AttributeError):
                        pass
                    
                    edited_chunks[text] = edited

                # 5. Replace original chunks with edited ones
                new_chunks = [edited_chunks.get(chunk, chunk) for chunk in chunks]

            # 6. Reconstruct and return the lesson
            return '\n\n'.join(new_chunks)
        except Exception as e:
            logger.error(f"Error in fallback RAG editing: {str(e)}", exc_info=True)
            return self._edit_lesson_simple(lesson_text, user_prompt)
    
    def _edit_lesson_simple(self, lesson_text: str, user_prompt: str) -> str:
        """Simple editing without RAG"""
        try:
            edit_prompt = f"""You are an expert teacher assistant. Please improve the following lesson content based on the user's request.

Lesson Content:
{lesson_text}

User Request:
{user_prompt}

Please provide an improved version of the lesson content that addresses the user's request. Return only the improved content."""

            response = self.llm.invoke(edit_prompt)
            if hasattr(response, 'content'):
                return response.content
            else:
                return str(response)
        except Exception as e:
            logger.error(f"Error in simple lesson editing: {str(e)}")
            return lesson_text

    def create_ppt(self, lesson_data: dict) -> bytes:
        """Generate a basic PPTX file from the lesson structure using python-pptx"""
        try:
            logger.info(f"Creating PowerPoint for lesson: {lesson_data.get('title', 'Unknown')}")
            logger.info(f"Lesson data keys: {list(lesson_data.keys())}")
            logger.info(f"Content length: {len(lesson_data.get('content', ''))}")
            
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = lesson_data.get('title', 'Lesson')
            slide.placeholders[1].text = lesson_data.get('summary', '')
            
            # Learning Objectives
            if lesson_data.get('learning_objectives'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Learning Objectives'
                body = slide.shapes.placeholders[1].text_frame
                for obj in lesson_data['learning_objectives']:
                    body.add_paragraph().text = str(obj)
            
            # Helper function to split content intelligently
            def split_content_into_chunks(content: str, max_length: int = 800) -> List[str]:
                """Split content into manageable chunks respecting paragraph boundaries"""
                if len(content) <= max_length:
                    return [content]
                
                chunks = []
                paragraphs = content.split('\n\n')
                current_chunk = ""
                
                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue
                    
                    # If adding this paragraph would exceed the limit
                    if current_chunk and len(current_chunk) + len(para) + 2 > max_length:
                        chunks.append(current_chunk)
                        current_chunk = para
                    else:
                        if current_chunk:
                            current_chunk += "\n\n" + para
                        else:
                            current_chunk = para
                    
                    # If a single paragraph is too long, split it by sentences
                    if len(para) > max_length:
                        if current_chunk:
                            chunks.append(current_chunk)
                            current_chunk = ""
                        # Split by sentences
                        sentences = para.replace('. ', '.\n').split('\n')
                        temp_chunk = ""
                        for sent in sentences:
                            if temp_chunk and len(temp_chunk) + len(sent) + 1 > max_length:
                                chunks.append(temp_chunk)
                                temp_chunk = sent
                            else:
                                temp_chunk += "\n" + sent if temp_chunk else sent
                        current_chunk = temp_chunk
                
                if current_chunk:
                    chunks.append(current_chunk)
                
                return chunks
            
            # Sections
            sections = lesson_data.get('sections', [])
            if sections:
                logger.info(f"Creating slides for {len(sections)} sections")
                for section in sections:
                    content = section.get('content', '')
                    heading = section.get('heading', 'Section')
                    
                    # Split content into manageable chunks
                    content_chunks = split_content_into_chunks(content)
                    
                    for chunk_idx, chunk in enumerate(content_chunks):
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        
                        # Set title - append part number if multiple slides
                        if len(content_chunks) > 1:
                            slide.shapes.title.text = f"{heading} (Part {chunk_idx + 1}/{len(content_chunks)})"
                        else:
                            slide.shapes.title.text = heading
                        
                        body = slide.shapes.placeholders[1].text_frame
                        body.text = chunk
                        
                        # Configure text formatting for readability
                        body.word_wrap = True
                        for paragraph in body.paragraphs:
                            paragraph.space_after = Pt(6)
                            paragraph.font.size = Pt(12)
            else:
                # If no sections, create a content slide with the main content
                if lesson_data.get('content'):
                    logger.info("Creating content slides with main lesson content")
                    content = lesson_data['content']
                    
                    # Split content into manageable chunks
                    content_chunks = split_content_into_chunks(content)
                    
                    for chunk_idx, chunk in enumerate(content_chunks):
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        
                        if len(content_chunks) > 1:
                            slide.shapes.title.text = f'Lesson Content (Part {chunk_idx + 1}/{len(content_chunks)})'
                        else:
                            slide.shapes.title.text = 'Lesson Content'
                        
                        body = slide.shapes.placeholders[1].text_frame
                        body.text = chunk
                        
                        # Configure text formatting for readability
                        body.word_wrap = True
                        for paragraph in body.paragraphs:
                            paragraph.space_after = Pt(6)
                            paragraph.font.size = Pt(12)
            
            # Key Concepts
            if lesson_data.get('key_concepts'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Key Concepts'
                body = slide.shapes.placeholders[1].text_frame
                for kc in lesson_data['key_concepts']:
                    body.add_paragraph().text = str(kc)
            
            # Background Prerequisites
            if lesson_data.get('background_prerequisites'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Background Prerequisites'
                body = slide.shapes.placeholders[1].text_frame
                for prereq in lesson_data['background_prerequisites']:
                    body.add_paragraph().text = str(prereq)
            
            # Creative Activities
            if lesson_data.get('creative_activities'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Creative Activities'
                body = slide.shapes.placeholders[1].text_frame
                for act in lesson_data['creative_activities']:
                    body.add_paragraph().text = f"{act.get('name', '')}: {act.get('description', '')} ({act.get('duration', '')})"
                    if act.get('learning_purpose'):
                        body.add_paragraph().text = f"Purpose: {act.get('learning_purpose', '')}"
            
            # STEM Equations
            if lesson_data.get('stem_equations'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'STEM Equations'
                body = slide.shapes.placeholders[1].text_frame
                for eq_data in lesson_data['stem_equations']:
                    if eq_data.get('equation'):
                        body.add_paragraph().text = f"Equation: {eq_data.get('equation', '')}"
                    if eq_data.get('complete_equation_significance'):
                        body.add_paragraph().text = f"Significance: {eq_data.get('complete_equation_significance', '')}"
            
            # Assessment Quiz
            if lesson_data.get('assessment_quiz'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Assessment Quiz'
                body = slide.shapes.placeholders[1].text_frame
                for q in lesson_data['assessment_quiz']:
                    body.add_paragraph().text = f"Q: {q.get('question', '')}"
                    for i, opt in enumerate(q.get('options', [])):
                        body.add_paragraph().text = f"{chr(65+i)}. {opt}"
                    body.add_paragraph().text = f"Answer: {q.get('answer', '')}"
            
            # Teacher Notes
            if lesson_data.get('teacher_notes'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Teacher Notes'
                body = slide.shapes.placeholders[1].text_frame
                for note in lesson_data['teacher_notes']:
                    body.add_paragraph().text = str(note)
            
            # If we only have a title slide, add a content slide
            if len(prs.slides) == 1 and lesson_data.get('content'):
                logger.info("Adding content slide as only title slide exists")
                content = lesson_data['content']
                
                # Split content into manageable chunks
                content_chunks = split_content_into_chunks(content)
                
                for chunk_idx, chunk in enumerate(content_chunks):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    
                    if len(content_chunks) > 1:
                        slide.shapes.title.text = f'Lesson Content (Part {chunk_idx + 1}/{len(content_chunks)})'
                    else:
                        slide.shapes.title.text = 'Lesson Content'
                    
                    body = slide.shapes.placeholders[1].text_frame
                    body.text = chunk
                    
                    # Configure text formatting for readability
                    body.word_wrap = True
                    for paragraph in body.paragraphs:
                        paragraph.space_after = Pt(6)
                        paragraph.font.size = Pt(12)
            
            logger.info(f"Created PowerPoint with {len(prs.slides)} slides")
            
            from io import BytesIO
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            ppt_bytes = buffer.getvalue()
            logger.info(f"PowerPoint generated successfully, size: {len(ppt_bytes)} bytes")
            return ppt_bytes
        except Exception as e:
            logger.error(f"Error creating PPTX: {str(e)}", exc_info=True)
            return b''

    def _create_docx_from_text(self, lesson_text: str, lesson_details: Optional[Dict[str, str]] = None) -> bytes:
        """Create DOCX from plain text lesson response"""
        try:
            doc = DocxDocument()
            
            # Add title
            title = lesson_details.get('lesson_title', 'Generated Lesson') if lesson_details else 'Generated Lesson'
            doc.add_heading(title, level=1)
            
            # Add lesson content
            doc.add_paragraph(lesson_text)
            
            # Save to bytes buffer
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error creating DOCX from text: {str(e)}")
            # Return a simple DOCX with error message
            doc = DocxDocument()
            doc.add_heading("Lesson Generation", level=1)
            doc.add_paragraph("A lesson has been generated from your content.")
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()

    def _create_docx(self, lesson_data: Dict[str, Any]) -> bytes:
        """Convert structured lesson to DOCX format with improved formatting"""
        try:
            doc = DocxDocument()
            
            # Add title with formatting
            title = doc.add_heading(level=1)
            title_run = title.add_run(self._sanitize_heading(lesson_data.get("title", "Generated Lesson")))
            title_run.bold = True
            
            # Add summary section
            if lesson_data.get("summary"):
                doc.add_heading(self._sanitize_heading("Summary"), level=2)
                summary = doc.add_paragraph(lesson_data["summary"])
                summary.paragraph_format.space_after = Inches(0.1)
            
            # Add learning objectives with bullet points
            if lesson_data.get("learning_objectives"):
                doc.add_heading(self._sanitize_heading("Learning Objectives"), level=2)
                for objective in lesson_data["learning_objectives"]:
                    p = doc.add_paragraph(style='ListBullet')
                    p.add_run(str(objective))
                doc.add_paragraph()
            
            # Add sections with proper spacing
            if lesson_data.get("sections"):
                for section in lesson_data["sections"]:
                    doc.add_heading(self._sanitize_heading(section.get("heading", "Section")), level=2)
                    content = doc.add_paragraph(section.get("content", ""))
                    content.paragraph_format.space_after = Inches(0.1)
                    doc.add_paragraph()
            
            # Add key concepts
            if lesson_data.get("key_concepts"):
                doc.add_heading(self._sanitize_heading("Key Concepts"), level=2)
                for concept in lesson_data["key_concepts"]:
                    p = doc.add_paragraph(style='ListBullet')
                    p.add_run(str(concept))
                doc.add_paragraph()
            
            # Add background prerequisites
            if lesson_data.get("background_prerequisites"):
                doc.add_heading(self._sanitize_heading("Background Prerequisites"), level=2)
                for prereq in lesson_data["background_prerequisites"]:
                    p = doc.add_paragraph(style='ListBullet')
                    p.add_run(str(prereq))
                doc.add_paragraph()
            
            # Add creative activities with clear formatting
            if lesson_data.get("creative_activities"):
                doc.add_heading(self._sanitize_heading("Creative Activities"), level=2)
                for i, activity in enumerate(lesson_data["creative_activities"], 1):
                    activity_title = doc.add_heading(level=3)
                    activity_title.add_run(self._sanitize_heading(f"Activity {i}: {activity.get('name', 'Unnamed Activity')}")).bold = True
                    
                    desc = doc.add_paragraph()
                    desc.add_run("Description: ").bold = True
                    desc.add_run(str(activity.get('description', '')))
                    
                    if activity.get('duration'):
                        duration = doc.add_paragraph()
                        duration.add_run("Duration: ").bold = True
                        duration.add_run(str(activity['duration']))
                    
                    if activity.get('learning_purpose'):
                        purpose = doc.add_paragraph()
                        purpose.add_run("Learning Purpose: ").bold = True
                        purpose.add_run(str(activity['learning_purpose']))
                    
                    doc.add_paragraph()
            
            # Add STEM equations if present
            if lesson_data.get("stem_equations"):
                doc.add_heading(self._sanitize_heading("STEM Equations"), level=2)
                for i, equation_data in enumerate(lesson_data["stem_equations"], 1):
                    if equation_data.get("equation"):
                        eq_title = doc.add_heading(level=3)
                        eq_title.add_run(self._sanitize_heading(f"Equation {i}")).bold = True
                        
                        eq_para = doc.add_paragraph()
                        eq_para.add_run("Equation: ").bold = True
                        eq_para.add_run(str(equation_data['equation']))
                        
                        if equation_data.get("term_explanations"):
                            terms = doc.add_paragraph()
                            terms.add_run("Term Explanations: ").bold = True
                            for term in equation_data["term_explanations"]:
                                term_p = doc.add_paragraph(style='ListBullet')
                                term_p.add_run(str(term))
                        
                        if equation_data.get("mathematical_operations"):
                            ops = doc.add_paragraph()
                            ops.add_run("Mathematical Operations: ").bold = True
                            ops.add_run(str(equation_data['mathematical_operations']))
                        
                        if equation_data.get("complete_equation_significance"):
                            sig = doc.add_paragraph()
                            sig.add_run("Complete Equation Significance: ").bold = True
                            sig.add_run(str(equation_data['complete_equation_significance']))
                        
                        doc.add_paragraph()
            
            # Add assessment quiz with clear question/answer formatting
            if lesson_data.get("assessment_quiz"):
                doc.add_heading(self._sanitize_heading("Assessment Quiz"), level=2)
                for i, question in enumerate(lesson_data["assessment_quiz"], 1):
                    q = doc.add_paragraph()
                    q.add_run(f"Question {i}: ").bold = True
                    q.add_run(str(question.get('question', '')))
                    
                    # Add options with letters
                    options = ['A', 'B', 'C', 'D']
                    question_options = question.get("options", [])
                    for opt, text in zip(options, question_options):
                        p = doc.add_paragraph(style='ListBullet')
                        p.add_run(f"{opt}. {str(text)}")
                    
                    # Add answer
                    ans = doc.add_paragraph()
                    ans.add_run("Correct Answer: ").bold = True
                    ans.add_run(str(question.get('answer', '')))
                    
                    if question.get('explanation'):
                        exp = doc.add_paragraph()
                        exp.add_run("Explanation: ").bold = True
                        exp.add_run(str(question['explanation']))
                    
                    doc.add_paragraph()
            
            # Add teacher notes
            if lesson_data.get("teacher_notes"):
                doc.add_heading(self._sanitize_heading("Teacher Notes"), level=2)
                for note in lesson_data["teacher_notes"]:
                    p = doc.add_paragraph(style='ListBullet')
                    p.add_run(str(note))
                doc.add_paragraph()
            
            # Save to bytes buffer
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error creating DOCX: {str(e)}")
            # Return a simple DOCX with error message
            doc = DocxDocument()
            doc.add_heading(self._sanitize_heading("Lesson Generation"), level=1)
            doc.add_paragraph("A lesson has been generated from your content.")
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()


