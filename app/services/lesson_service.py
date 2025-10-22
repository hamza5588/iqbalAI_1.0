<<<<<<< HEAD
"""
Main lesson service - now uses structured approach with separate teacher and student services
"""
=======
import os
>>>>>>> 342f1c22ac17d5654a67a312daa13bd9dfd46cdd
import logging
from typing import Any, Dict, Optional
from werkzeug.datastructures import FileStorage

from .lesson.teacher_service import TeacherLessonService
from .lesson.student_service import StudentLessonService

# Set up logging
logger = logging.getLogger(__name__)


class LessonService:
    """
    Main lesson service that delegates to specialized teacher and student services.
    This maintains backward compatibility while providing a structured approach.
    """
    def __init__(self, api_key: str):
        """Initialize the LessonService with API key from database."""
        self.api_key = api_key
        self.teacher_service = TeacherLessonService(api_key)
        self.student_service = StudentLessonService(api_key)

    # Delegate methods to appropriate services
    
    def allowed_file(self, filename: str) -> bool:
        """Check if file extension is supported"""
        return self.teacher_service.allowed_file(filename)

    def process_file(self, file: FileStorage, lesson_details: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
        """Process an uploaded file and return structured lesson content with DOCX bytes."""
        return self.teacher_service.process_file(file, lesson_details)




    def create_ppt(self, lesson_data: dict) -> bytes:
<<<<<<< HEAD
        """Generate a basic PPTX file from the lesson structure using python-pptx."""
        return self.teacher_service.create_ppt(lesson_data)

    def edit_lesson_with_prompt(self, lesson_text: str, user_prompt: str) -> str:
        """Use a FAISS vector database for semantic chunk retrieval and editing."""
        return self.teacher_service.edit_lesson_with_prompt(lesson_text, user_prompt)
=======
            
        """
        Generate a PowerPoint presentation where:
        - Slide 0: Title only
        - Following slides: All content from sections (summary + content combined)
        - Final slide: Learning objectives
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        import logging

        logger = logging.getLogger(__name__)

        try:
            prs = Presentation()
            logger.info(f"Creating PowerPoint for lesson: {lesson_data.get('title', 'Unknown')}")

            # --- SLIDE 0: Title Only ---
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = lesson_data.get('title', 'Lesson')

            # --- All Content Slides (Summary + Content from sections) ---
            sections = lesson_data.get('sections', [])
            
            if sections:
                for section in sections:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = section.get('heading', 'Content')

                    # Add textbox with content
                    left, top = Inches(1), Inches(1.8)
                    width, height = Inches(8.5), Inches(4.5)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.word_wrap = True
                    text_frame.vertical_anchor = 1  # Top alignment

                    p = text_frame.paragraphs[0]
                    p.text = section.get('content', '')
                    p.font.size = Pt(18)
                    p.line_spacing = 1.3
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
            else:
                # Fallback if no sections provided
                logger.warning("No sections provided for PPT generation")

            # --- Learning Objectives Slide ---
            if lesson_data.get('learning_objectives'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = 'Learning Objectives'
                
                left, top = Inches(1), Inches(1.8)
                width, height = Inches(8.5), Inches(4.5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                
                for obj in lesson_data['learning_objectives']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {obj}"
                    p.font.size = Pt(16)
                    p.line_spacing = 1.3
                    p.space_before = Pt(6)

            # --- Save Presentation ---
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            logger.info(f"✅ PPT created successfully with {len(prs.slides)} slides.")
            return buffer.getvalue()

        except Exception as e:
            logger.error(f"Error creating PPTX: {str(e)}", exc_info=True)
            return b''


    def answer_lesson_question(self, lesson_id, question):
        # Get lesson content
        lesson = LessonModel.get_lesson_by_id(lesson_id)
        if not lesson:
            # Handle demo lessons that don't exist in database
            if lesson_id in [1, 2, 3]:
                demo_lessons = {
                    1: {
                        'title': 'Introduction to Climate',
                        'content': 'Climate refers to the long-term patterns of temperature, humidity, wind, and precipitation in a particular region. Unlike weather, which describes short-term atmospheric conditions, climate represents the average weather conditions over a period of 30 years or more. Key concepts include climate vs weather, climate factors (latitude, altitude, ocean currents, wind patterns, landforms), and climate zones (tropical, temperate, polar, desert). Climate change is a significant topic, with human activities like burning fossil fuels contributing to increased greenhouse gas concentrations.'
                    },
                    2: {
                        'title': 'World History Overview',
                        'content': 'World history encompasses the study of human civilization from its earliest beginnings to the present day. Major periods include Ancient Civilizations (Mesopotamia, Egypt, Greece, Rome), Middle Ages (500-1500 CE), Early Modern Period (1500-1800), and Modern Period (1800-Present). Key themes include power and politics, economics, culture, technology, and environment. Understanding world history helps us make sense of current events and prepare for future challenges.'
                    },
                    3: {
                        'title': 'Physics Fundamentals',
                        'content': 'Physics is the study of matter, energy, and their interactions. Core areas include mechanics (Newton\'s laws, forces, energy, momentum), waves and sound (wave properties, sound waves), electricity and magnetism (electric charges, current, magnetism), and modern physics (quantum mechanics, relativity). Physics principles apply to engineering, medicine, technology, space exploration, and energy. Understanding physics helps us comprehend the natural world and develop technologies that improve our lives.'
                    }
                }
                demo_lesson = demo_lessons[lesson_id]
                lesson_title = demo_lesson['title']
                content = demo_lesson['content']
                logger.info(f"Using demo lesson '{lesson_title}' for ID: {lesson_id}")
            else:
                return {'error': 'Lesson not found'}
        else:
            content = lesson.get('content') or lesson.get('summary')
            if not content:
                return {'error': 'No lesson content available'}
            lesson_title = lesson.get('title', 'this lesson')
            logger.info(f"Answering question for lesson '{lesson_title}' (ID: {lesson_id})")
        
        # Use LLM to answer with lesson-specific context
        answer = self.llm_answer(content, question, lesson_title)

        # Canonicalize the question and log to FAQ using semantic matching
        try:
            canonical = self.canonicalize_question(lesson_id, question)
        except Exception:
            canonical = question
        LessonFAQ.log_question(lesson_id, canonical)
        return {'answer': answer, 'canonical_question': canonical}

    def llm_answer(self, lesson_content, question, lesson_title="this lesson"):
        # Check if the question references specific lines
        line_references = self._extract_line_references(question)
        
        # Use Groq LLM (already implemented in the project)
        if line_references:
            # Format lesson content with line numbers for context
            numbered_content = self._format_content_with_line_numbers(lesson_content)
            prompt = f"""
You are a helpful teacher for the lesson titled "{lesson_title}". The student is asking about specific lines in the lesson content. Use the following numbered lesson content to answer their question concisely and clearly.

Lesson Title: {lesson_title}
Numbered Lesson Content:
{numbered_content}

Student's Question: {question}

The student is specifically asking about line(s): {', '.join(map(str, line_references))}

Answer as a helpful teacher, focusing on the specific lines mentioned and providing clear explanations for those lines:
"""
        else:
            prompt = f"""
You are a helpful teacher for the lesson titled "{lesson_title}". Use the following lesson content to answer the student's question concisely and clearly. Make sure your response is specific to this lesson and its content.

Lesson Title: {lesson_title}
Lesson Content:
{lesson_content}

Student's Question: {question}

Answer as a helpful teacher, being specific to the "{lesson_title}" lesson content:
"""
        try:
            response = self.llm.invoke(prompt)
            if hasattr(response, 'content'):
                return response.content.strip()
            else:
                return str(response).strip()
        except Exception as e:
            return f"[Error from LLM: {e}]"

    def _extract_line_references(self, question: str) -> list:
        """Extract line references from a question like 'line 1', 'sentence 2', etc."""
        import re
        
        line_patterns = [
            r'line\s+(\d+)',
            r'sentence\s+(\d+)',
            r'paragraph\s+(\d+)',
            r'number\s+(\d+)',
            r'(\d+)(?:st|nd|rd|th)?\s+(?:line|sentence|paragraph)'
        ]
        
        references = []
        for pattern in line_patterns:
            matches = re.findall(pattern, question, re.IGNORECASE)
            for match in matches:
                line_number = int(match)
                if line_number > 0:
                    references.append(line_number)
        
        return references

    def _format_content_with_line_numbers(self, content: str) -> str:
        """Format lesson content with line numbers for AI context"""
        import re
        
        # Split content into paragraphs and sentences
        paragraphs = content.split('\n\n')
        numbered_content = ''
        line_number = 1
        
        for paragraph in paragraphs:
            if paragraph.strip():
                # Split paragraph into sentences
                sentences = re.split(r'[.!?]+', paragraph)
                sentences = [s.strip() for s in sentences if s.strip()]
                
                if len(sentences) > 1:
                    # Multiple sentences - number each sentence
                    for sentence in sentences:
                        if sentence:
                            numbered_content += f"{line_number}: {sentence}.\n"
                            line_number += 1
                else:
                    # Single sentence or short paragraph - number as one line
                    numbered_content += f"{line_number}: {paragraph.strip()}\n"
                    line_number += 1
                
                numbered_content += "\n"
        
        return numbered_content

    def canonicalize_question(self, lesson_id: int, question: str) -> str:
        """Return a canonical phrasing for the question using semantic similarity.

        - Fetch existing canonical questions for the lesson
        - If any are semantically similar above a threshold, reuse that canonical
        - Otherwise, create a concise canonical phrasing via the LLM
        """
        try:
            import sqlite3
            from langchain_community.embeddings import HuggingFaceEmbeddings

            conn = sqlite3.connect('instance/chatbot.db')
            c = conn.cursor()
            c.execute('''CREATE TABLE IF NOT EXISTS lesson_faq (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                lesson_id INTEGER,
                question TEXT,
                count INTEGER DEFAULT 1,
                canonical_question TEXT
            )''')
            c.execute('SELECT COALESCE(canonical_question, question) FROM lesson_faq WHERE lesson_id = ?', (lesson_id,))
            existing = [row[0] for row in c.fetchall() if row and row[0]]
            conn.close()

            if not existing:
                # No prior questions; generate a canonical phrasing
                return self._generate_canonical(question)

            # Embed existing and incoming question
            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            existing_vectors = embeddings.embed_documents(existing)
            query_vector = embeddings.embed_query(question)

            # Compute cosine similarity
            def cosine(a, b):
                import math
                dot = sum(x*y for x, y in zip(a, b))
                na = math.sqrt(sum(x*x for x in a))
                nb = math.sqrt(sum(y*y for y in b))
                return (dot / (na * nb)) if na and nb else 0.0

            sims = [cosine(vec, query_vector) for vec in existing_vectors]
            best_idx, best_sim = (max(enumerate(sims), key=lambda t: t[1]) if sims else (-1, 0.0))

            # Threshold tuned for paraphrase equivalence for MiniLM embeddings
            if best_sim >= 0.82:
                return existing[best_idx]

            # Otherwise, produce a concise canonical phrasing
            return self._generate_canonical(question)
        except Exception:
            # On any error, fall back to original question
            return question

    def _generate_canonical(self, question: str) -> str:
        """Use the LLM to produce a short canonical phrasing for a question."""
        try:
            prompt = (
                "Rewrite the following student question into a short, neutral, canonical phrasing (max 12 words).\n"
                "Keep the meaning identical. Do not include quotes or extra commentary.\n\n"
                f"Question: {question}\n\nCanonical:"
            )
            response = self.llm.invoke(prompt)
            text = response.content.strip() if hasattr(response, 'content') else str(response).strip()
            # Clean up extraneous punctuation/quotes
            text = text.strip('"\'\u201c\u201d ')
            return text or question
        except Exception:
            return question

    def get_lesson_faqs(self, lesson_id, limit=5):
        """Get frequently asked questions for a lesson"""
        try:
            # Get lesson content from database
            from app.models.models import LessonModel
            lesson = LessonModel.get_lesson_by_id(lesson_id)
            if not lesson:
                return []
            
            lesson_content = lesson['content']
            
            # Generate FAQs using LLM
            faq_prompt = f"""
            Based on the following lesson content, generate {limit} frequently asked questions that students might have.
            Focus on key concepts, common misconceptions, and important details.
            
            Lesson Content:
            {lesson_content}
            
            Return the questions in this JSON format:
            {{
                "faqs": [
                    {{
                        "question": "What is...?",
                        "answer": "The answer is..."
                    }}
                ]
            }}
            """
            
            response = self.llm.invoke(faq_prompt)
            response_text = response.content
            
            # Parse JSON response
            try:
                import json
                faq_data = json.loads(response_text)
                return faq_data.get('faqs', [])
            except json.JSONDecodeError:
                logger.error(f"Failed to parse FAQ JSON: {response_text}")
                return []
                
        except Exception as e:
            logger.error(f"Error generating FAQs: {str(e)}")
            return []
>>>>>>> 342f1c22ac17d5654a67a312daa13bd9dfd46cdd

    def improve_lesson_content(self, lesson_id: int, current_content: str, improvement_prompt: str = "") -> str:
        """Improve lesson content using AI based on user prompt"""
        return self.teacher_service.improve_lesson_content(lesson_id, current_content, improvement_prompt)

    def review_lesson_with_rag(self, lesson_content: str, user_prompt: str, filename: str = "") -> str:
        """Review lesson content using RAG to retrieve relevant information from vector database"""
        return self.teacher_service.review_lesson_with_rag(lesson_content, user_prompt, filename)

    # Student-focused methods
    def answer_lesson_question(self, lesson_id: int, question: str) -> Dict[str, str]:
        """Answer a student's question about a specific lesson"""
        return self.student_service.answer_lesson_question(lesson_id, question)

    def get_lesson_faqs(self, lesson_id: int, limit: int = 5) -> list:
        """Get frequently asked questions for a lesson"""
        return self.student_service.get_lesson_faqs(lesson_id, limit)

    def get_lesson_summary(self, lesson_id: int) -> Dict[str, str]:
        """Get a summary of the lesson for students"""
        return self.student_service.get_lesson_summary(lesson_id)

    def get_lesson_key_points(self, lesson_id: int) -> list:
        """Extract key learning points from a lesson"""
        return self.student_service.get_lesson_key_points(lesson_id)

    # Legacy methods for backward compatibility
    def llm_answer(self, lesson_content: str, question: str, lesson_title: str = "this lesson") -> str:
        """Generate an answer using the LLM with lesson-specific context"""
        return self.student_service.llm_answer(lesson_content, question, lesson_title)

    def canonicalize_question(self, lesson_id: int, question: str) -> str:
        """Return a canonical phrasing for the question using semantic similarity"""
        return self.student_service.canonicalize_question(lesson_id, question)